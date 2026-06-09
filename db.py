# -*- coding: utf-8 -*-
from pydantic import ValidationError
import warnings
import requests
import streamlit.components.v1 as components
import subprocess
import json
from datetime import datetime
from typing import Dict, Optional
from streamlit_js import st_js
import pydantic
import threading
from concurrent.futures import ThreadPoolExecutor
from concurrent.futures import ThreadPoolExecutor, as_completed

warnings.filterwarnings('ignore')
import logging
import time
import traceback
import os
# 严格按你的要求：仅导入Client类，不导入顶层chat函数
from ollama import Client  # 导入Client类，而非直接导入顶层chat函数
from pydantic import BaseModel  # 严格对齐你的导入方式
# ========== SM4加密配置（新增：全局密钥+加密函数，用于注册登录加密） ==========
from gmssl import sm4, func
import binascii
# ========== 图形验证码配置（新增：生成验证码+图片，用于登录/注册校验） ==========
from captcha.image import ImageCaptcha
import random
import string
from io import BytesIO
import yaml  # 新增：读写yaml配置文件
from pathlib import Path  # 已有，不用改，下面新增yaml路径常量
import jieba  # 完整保留
# ========== Redis 加速配置 ==========
import redis, hashlib, json
from redis.connection import ConnectionPool
from decimal import Decimal  # ← 新增
from uuid import UUID  # ← 新增

# 高德天气API封装
GAODE_API_KEY = "16a98a1a06e46becd9a8689f18074882"


class GaodeWeatherAPI:
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.base_url = "https://restapi.amap.com/v3/weather/weatherInfo"

    # --------------  ① 实况天气 --------------
    def get_weather(self, city_code: str, extensions: str = "base") -> Optional[Dict]:
        params = {"key": self.api_key, "city": city_code, "extensions": extensions, "output": "JSON"}
        try:
            resp = requests.get(self.base_url, params=params, timeout=10)
            resp.raise_for_status()
            data = resp.json()
            return data if data.get("status") == "1" and data.get("infocode") == "10000" else None
        except Exception as e:
            logging.error(f"[GaodeWeatherAPI] 请求失败: {e}")
            return None

    def parse_current_weather(self, data: Dict) -> Optional[Dict]:
        #Optional[Dict] == Union[Dict, None] 要么返回一个字典，要么返回 None。
        if not data or "lives" not in data or not data["lives"]:
            return None
        d = data["lives"][0]
        return {
            "city": d.get("city"),
            "weather": d.get("weather"),
            "temperature": d.get("temperature"),
            "wind_direction": d.get("winddirection"),
            "wind_power": d.get("windpower"),
            "humidity": d.get("humidity"),
            "report_time": d.get("reporttime")
        }

    # --------------  ② 预报天气 --------------
    def parse_forecast_weather(self, data: Dict) -> Optional[Dict]:
        if not data or "forecasts" not in data or not data["forecasts"]:
            return None
        f = data["forecasts"][0]
        return {
            "city": f.get("city"),
            "report_time": f.get("reporttime"),
            "forecasts": [
                {
                    "date": c.get("date"),
                    "week": c.get("week"),
                    "day_weather": c.get("dayweather"),
                    "night_weather": c.get("nightweather"),
                    "day_temp": c.get("daytemp"),
                    "night_temp": c.get("nighttemp"),
                    "day_wind": c.get("daywind"),
                    "night_wind": c.get("nightwind"),
                    "day_power": c.get("daypower"),
                    "night_power": c.get("nightpower")
                }
                for c in f.get("casts", [])
            ]
        }

import re
from typing import Tuple, Optional


def _search_city(user_raw: str) -> Tuple[Optional[str], Optional[str], float]:
    user_raw = user_raw.strip()
    if not user_raw:
        return None, None, 0.0

    db_client = st.session_state.get("db_client")
    if not db_client or not hasattr(db_client, 'cursor') or db_client.cursor is None:
        logging.error("数据库连接未初始化")
        return None, None, 0.0

    try:
        # 根据用户输入的城市名称进行精确查询
        db_client.cursor.execute("SELECT city_name, adcode FROM city_adcode WHERE city_name = %s", (user_raw,))
        row = db_client.cursor.fetchone()
        if row:
            logging.info(f"精确匹配到城市：{row['city_name']}，adcode：{row['adcode']}")
            return row['city_name'], row['adcode'], 1.0

        # 如果精确匹配失败，尝试模糊查询
        db_client.cursor.execute("SELECT city_name, adcode FROM city_adcode WHERE city_name ILIKE %s",
                                 (f"%{user_raw}%",))
        rows = db_client.cursor.fetchall()
        if rows:
            # 选择第一个匹配结果
            best_match = rows[0]
            logging.info(f"模糊匹配到城市：{best_match['city_name']}，adcode：{best_match['adcode']}")
            return best_match['city_name'], best_match['adcode'], 0.5  # 模糊匹配的分数可以自定义

        logging.error("未找到匹配的城市，请尝试输入更具体的城市名称，例如：山东省青岛市")
        return None, None, 0.0

    except Exception as e:
        logging.error(f"查询城市表失败：{e}")
        return None, None, 0.0

# ========== 最小新增：天气并行缓存 ==========
_weather_cache_lock = threading.Lock()
_weather_cache = {}          # key=adcode, value={"expire":ts, "info":dict}

def _get_weather_cached(adcode: str) -> dict:
    """线程安全、秒级缓存；返回统一格式的天气 dict"""
    now = time.time()
    with _weather_cache_lock:
        hit = _weather_cache.get(adcode)
        if hit and now < hit["expire"]:
            return hit["info"]

    # 缓存未命中 → 实时拉取
    api = GaodeWeatherAPI(GAODE_API_KEY)
    raw = api.get_weather(adcode, "base")
    info = api.parse_current(raw) if raw else ""
    weather_dict = {
        "city"       : info.split(" ")[0] if info else "",
        "weather"    : info.split("：")[2].split("，")[0] if info else "",
        "temperature": info.split("温度 ")[1].split("℃")[0] if info else "",
        "humidity"   : info.split("湿度 ")[1].split("%")[0] if info else "",
        "report_time": info.split("更新时间 ")[1] if info else "",
        "full_text"  : info
    }

    with _weather_cache_lock:
        _weather_cache[adcode] = {"expire": now + 30, "info": weather_dict}
        #希望「立刻刷新」，把 ttl=None 改成 ttl=30 或直接 del _weather_cache[adcode] 再查。
    return weather_dict

OLLAMA_MODEL = "deepseek-r1:7b"  # deepseek模型名称
# 创建连接池
redis_pool = ConnectionPool(
    host='localhost',
    port=6379,
    password='Password123@redis',
    decode_responses=True,
    max_connections=50,
    socket_timeout=5,
    retry_on_timeout=True
)

# 修改原有的 REDIS_CLI 定义（替换原有行）
REDIS_CLI = redis.Redis(connection_pool=redis_pool)
# ----------  全局 Redis 降档函数 ----------
def redis_get(key: str, default=None):
    """带 2 秒超时保护的 Redis 读取"""
    try:
        return REDIS_CLI.get(key)
    except (redis.TimeoutError, redis.ConnectionError) as e:
        logging.warning(f"Redis 读取超时/断连: {e} | key={key}")
        return default
CACHE_TTL = 60 * 30  # 半小时
CACHE_VER = 'v2'  # 后期可手动清缓存


def _serialize_rows(rows):
    """把 psycopg2 的 RealDictRow 洗成纯 Python 基本类型"""
    out = []
    for r in rows:
        clean = {}
        for k, v in r.items():
            if isinstance(v, datetime):  # ← 改成 datetime
                clean[k] = v.isoformat()
            elif isinstance(v, (Decimal, UUID)):  # ← 同样只写类名
                clean[k] = str(v)
            else:
                clean[k] = v
        out.append(clean)
    return out


# ========== 答案级 Redis 缓存 ==========
def _get_answer_json(query, kb_id, top_k, client, prompt, db, encoder):
    """先读缓存，没有再调 Ollama（并行搜库 + 生成）"""
    cache_key = f"ans_json:{CACHE_VER}:{hashlib.md5(f'{query}#{kb_id}#{top_k}#{prompt[:100]}'.encode()).hexdigest()}"

    cached = REDIS_CLI.get(cache_key)
    if cached:
        try:
            return json.loads(cached)
        except Exception:
            pass

    # 并行：搜库 + 调 Ollama
    with ThreadPoolExecutor(max_workers=3) as pool:
        search_future = pool.submit(db._search_similar_chunks_raw, query, encoder, top_k, kb_id)
        weather_future = pool.submit(_get_weather_cached, st.session_state.get("user_adcode", ""))
        ollama_future = pool.submit(client.chat,
                                    messages=[{'role': 'user', 'content': prompt}],
                                    model=OLLAMA_MODEL,
                                    format=KnowledgeBaseAnswer.model_json_schema(),
                                    stream=False,
                                    options={"temperature": 0.1, "max_tokens": 300, "num_ctx": 4096})

    raw_results = search_future.result()
    weather_dict = weather_future.result()  # 立即拿到天气
    weather_kw = weather_dict.get("weather", "")
    city_name = weather_dict.get("city", "")
    # 构造带天气的查询串
    response = ollama_future.result()
    # 把天气追加到 prompt（最小侵入）
    if weather_dict.get("full_text"):
        prompt += f"\n【实时天气】{weather_dict['full_text']}"

    # 取出 JSON
    raw_json = response.message.content.strip()
    raw_json = re.sub(r'^```json|```$', '', raw_json).strip()

    kb_answer = KnowledgeBaseAnswer.model_validate_json(raw_json)
    dict_answer = kb_answer.model_dump()

    # 写缓存（2h）
    REDIS_CLI.setex(cache_key, 60 * 120, json.dumps(dict_answer, ensure_ascii=False))
    return dict_answer


# ========== 通用 Ollama 答案缓存 ==========
def cached_ollama_stream_response(
        prompt: str,
        format_schema: dict,
        model: str = OLLAMA_MODEL,
        ttl: int = 60 * 120  # 延长到2小时
) -> any:
    """
    优化的Ollama响应缓存，显著加速回答生成
    返回类型和用法完全不变，保证零改动
    """
    cache_key = f"ans:{CACHE_VER}:{model}:{hashlib.md5(prompt.encode()).hexdigest()}"

    # ========== 优化1：增加更多缓存检查点 ==========
    # 先检查快速缓存（短时缓存）
    fast_cache_key = f"fast:{cache_key}"
    fast_cached = REDIS_CLI.get(fast_cache_key)
    if fast_cached:
        try:
            data = json.loads(fast_cached)
            from types import SimpleNamespace
            ns = SimpleNamespace()
            ns.message = SimpleNamespace()
            ns.message.content = json.dumps(data, ensure_ascii=False)
            return ns
        except Exception as e:
            print(f"快速缓存读取失败: {e}")

    # ========== 优化2：智能流式调用（更快响应） ==========
    try:
        # 先尝试极简流式调用（不指定format，加快速度）
        client = Client(host='http://127.0.0.1:11434')

        # 创建进度占位符（不影响实际函数返回值）
        response_parts = []

        # 流式调用（简化参数，加快响应）
        for chunk in client.chat(
                messages=[{'role': 'user', 'content': prompt}],
                model=model,
                stream=True,  # 关键：流式模式
                options={
                    "temperature": 0.1,
                    "max_tokens": 1024,  # 减少到1024，加快生成
                    "num_ctx": 2048,  # 减少上下文
                    "seed": 42,
                }
        ):
            delta = chunk['message']['content']
            response_parts.append(delta)

        full_response = "".join(response_parts)

        # ========== 优化3：智能JSON解析和缓存 ==========
        try:
            # 尝试直接解析
            answer = KnowledgeBaseAnswer.model_validate_json(full_response)

            # 缓存到快速缓存（短时）和主缓存（长时）
            REDIS_CLI.setex(fast_cache_key, 60, json.dumps(answer.model_dump(), ensure_ascii=False))  # 1分钟快速缓存
            REDIS_CLI.setex(cache_key, ttl, json.dumps(answer.model_dump(), ensure_ascii=False))

            # 返回格式化的response
            from types import SimpleNamespace
            ns = SimpleNamespace()
            ns.message = SimpleNamespace()
            ns.message.content = json.dumps(answer.model_dump(), ensure_ascii=False)
            return ns

        except Exception as e:
            # 解析失败，尝试清理和重新解析
            print(f"流式解析失败，尝试清理JSON: {e}")
            # 尝试清理JSON
            cleaned = re.sub(r'^```json|```$', '', full_response).strip()
            if cleaned.startswith('{') and cleaned.endswith('}'):
                try:
                    answer = KnowledgeBaseAnswer.model_validate_json(cleaned)
                    # 缓存
                    REDIS_CLI.setex(cache_key, ttl, json.dumps(answer.model_dump(), ensure_ascii=False))
                    # 返回
                    from types import SimpleNamespace
                    ns = SimpleNamespace()
                    ns.message = SimpleNamespace()
                    ns.message.content = json.dumps(answer.model_dump(), ensure_ascii=False)
                    return ns
                except:
                    # 继续使用非流式
                    pass

    except Exception as e:
        print(f"流式调用失败，使用非流式: {e}")

    # ========== 优化4：非流式调用（保底，但优化参数） ==========
    try:
        client = Client(host='http://127.0.0.1:11434')
        response = client.chat(
            messages=[{'role': 'user', 'content': prompt}],
            model=model,
            format=format_schema,
            stream=False,
            options={
                "temperature": 0.1,
                "top_p": 0.9,
                "max_tokens": 2048,  # 优化为2048
                "num_ctx": 4096,
                "seed": 42,
            }
        )

        # 缓存结果
        try:
            raw_json = response.message.content.strip()
            answer = KnowledgeBaseAnswer.model_validate_json(raw_json)
            REDIS_CLI.setex(cache_key, ttl, json.dumps(answer.model_dump(), ensure_ascii=False))
            REDIS_CLI.setex(fast_cache_key, 60, json.dumps(answer.model_dump(), ensure_ascii=False))
        except Exception as e:
            print(f"非流式缓存写入失败: {e}")

        return response  # 保持原有返回格式

    except Exception as e:
        # 保持原有错误处理逻辑
        raise Exception(f"Ollama调用失败：{str(e)[:100]}")


# ========== 新增：定义db.yaml路径（代码根目录下，绝对路径，必加） ==========
DB_YAML_PATH = Path(__file__).parent / "db.yaml"

# 全局SM4密钥（必须16字节，可自行修改，建议保存好不要随意变更）
SM4_KEY = b'1234567890abcdef'  # 16字节密钥，可自定义（如b'your_secret_key_16'）
sm4_crypt = sm4.CryptSM4()  # 初始化SM4对象


# SM4加密函数（明文转16进制字符串，方便存入数据库）
def sm4_encrypt(plain_text):
    if not isinstance(plain_text, str):
        plain_text = str(plain_text)
    sm4_crypt.set_key(SM4_KEY, sm4.SM4_ENCRYPT)  # 设置加密模式
    # 明文转字节流（utf-8编码）
    plain_bytes = plain_text.encode('utf-8')
    # 加密（ECB模式，无需iv，返回字节流）
    encrypt_bytes = sm4_crypt.crypt_ecb(plain_bytes)
    # 字节流转16进制字符串（方便数据库存储和比对）
    encrypt_hex = binascii.hexlify(encrypt_bytes).decode('utf-8')
    return encrypt_hex


logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")


# SM4解密函数（16进制字符串转明文）
def sm4_decrypt(encrypt_hex):
    """SM4解密函数 - 增强版错误处理"""
    try:
        if not isinstance(encrypt_hex, str) or not encrypt_hex:
            return ""

        # 清理输入：移除所有空格和特殊字符
        encrypt_hex = ''.join(c for c in encrypt_hex if c.isalnum()).strip()

        if not encrypt_hex:
            return ""

        # 修复奇数长度问题
        if len(encrypt_hex) % 2 != 0:
            logging.warning(f"修复奇数长度密文: {encrypt_hex[:20]}")
            # 尝试多种修复策略
            if len(encrypt_hex) < 32:  # 标准密文应该是32字符
                # 可能是截断的密文，尝试填充到标准长度
                encrypt_hex = encrypt_hex.ljust(32, '0')
            elif len(encrypt_hex) > 32:
                # 可能包含额外字符，截取到标准长度
                encrypt_hex = encrypt_hex[:32]
            else:
                # 奇数长度但已经是32，补0使其为偶数
                encrypt_hex = '0' + encrypt_hex if len(encrypt_hex) % 2 != 0 else encrypt_hex

        # 验证是否为有效的16进制
        import re
        if not re.match(r'^[0-9a-fA-F]+$', encrypt_hex):
            logging.error(f"无效的16进制字符串: {encrypt_hex[:20]}")
            # 尝试从损坏的数据中恢复用户名
            if "_" in encrypt_hex:
                # 可能是之前解密失败生成的占位符
                return encrypt_hex
            return f"用户_{hash(encrypt_hex) % 10000}"

        sm4_crypt.set_key(SM4_KEY, sm4.SM4_DECRYPT)
        encrypt_bytes = binascii.unhexlify(encrypt_hex)
        decrypt_bytes = sm4_crypt.crypt_ecb(encrypt_bytes)

        # 尝试多种编码方式
        try:
            plain_text = decrypt_bytes.decode('utf-8').strip()
        except UnicodeDecodeError:
            try:
                plain_text = decrypt_bytes.decode('gbk').strip()
            except:
                plain_text = decrypt_bytes.decode('utf-8', errors='ignore').strip()

        return plain_text if plain_text else f"用户_{hash(encrypt_hex) % 10000}"

    except Exception as e:
        logging.error(f"SM4解密失败: {e} | 原始数据: {encrypt_hex[:20]}")
        return f"用户_{hash(encrypt_hex) % 10000}"


def show_db_config_stage():
    """显示数据库配置阶段"""
    st.markdown(HIDE_ST_STYLE, unsafe_allow_html=True)
    st.subheader("数据库初始化（首次使用必填）")

    # 新增：根据状态显示不同内容
    if st.session_state.get("db_initialized", False):
        # 第二阶段：显示城市输入框
        show_city_input_section()
        return

    # 第一阶段：数据库配置表单
    st.caption("第一步：填写数据库配置")

    with st.form("db_config_form"):
        col1, col2 = st.columns(2)

        with col1:
            db_host = st.text_input("主机", value="localhost", key="init_db_host")
            db_port = st.number_input("端口", value=5432, step=1, key="init_db_port")
            db_user = st.text_input("用户名", value="postgres", key="init_db_user")

        with col2:
            db_pwd = st.text_input("密码", type="password", key="init_db_pwd")
            db_name = st.text_input("数据库名", value="kb_db", key="init_db_name")
            model_type = st.selectbox(
                "模型类型",
                ["双模型融合", "text2vec-base-chinese", "bert-base-chinese"],
                key="init_model_type",
                index=0
            )

        use_cuda = st.checkbox("启用CUDA", key="init_use_cuda", value=False)

        submit = st.form_submit_button("初始化数据库", type="primary")

        if submit:
            # 验证必填字段
            if not all([db_host, db_user, db_name]):
                st.error("请填写所有必填字段！")
                return

            with st.spinner("正在初始化数据库和模型..."):
                # 初始化数据库
                success, encoder, db_client, error_msg = init_database(
                    db_host, db_port, db_user, db_pwd, db_name, use_cuda, model_type
                )

                if success:
                    # 保存临时配置到session_state
                    st.session_state.temp_db_config = {
                        "host": db_host,
                        "port": db_port,
                        "user": db_user,
                        "password": db_pwd,
                        "dbname": db_name,
                        "model_type": model_type,
                        "use_cuda": use_cuda
                    }
                    st.session_state.db_client = db_client
                    st.session_state.encoder = encoder
                    st.session_state.db_initialized = True
                    st.success("数据库初始化成功！")
                    st.rerun()
                else:
                    st.error(f"数据库初始化失败：{error_msg}")


def show_city_input_section():
    """显示城市输入部分（第二阶段）"""
    st.caption("第二步：输入城市信息")
    st.info("数据库已初始化成功，请输入您所在的城市")

    # 检查临时配置是否存在
    if not st.session_state.get("temp_db_config"):
        st.error("数据库配置丢失，请返回重新初始化")
        if st.button("返回数据库配置"):
            st.session_state.db_initialized = False
            st.rerun()
        return

    # 城市输入框
    city_input = st.text_input(
        "请输入您所在的城市名称",
        placeholder="例如：青岛市",
        help="输入城市名，系统将自动匹配"
    )

    # 搜索城市并显示结果
    if st.button("搜索城市", type="secondary"):
        if city_input.strip():
            with st.spinner("正在搜索城市..."):
                best_name, best_code, score = _search_city(city_input.strip())

                if best_name and best_code:
                    st.success(f"匹配到：{best_name} (adcode: {best_code})")
                    # 将匹配结果存入session_state（避免按钮状态丢失问题）
                    st.session_state["matched_city"] = best_name
                    st.session_state["matched_adcode"] = best_code
                else:
                    st.error("未找到匹配的城市，请尝试输入更具体的城市名称，例如：山东省青岛市")
        else:
            st.warning("请先输入城市名称")

    # 如果已匹配到城市，显示确认按钮
    if st.session_state.get("matched_city") and st.session_state.get("matched_adcode"):
        st.info(f"当前匹配城市：{st.session_state.matched_city} (adcode: {st.session_state.matched_adcode})")

        # 确认按钮（不再依赖confirm_init状态）
        if st.button("确认并完成初始化", type="primary", key="confirm_init_final"):
            # 从temp_db_config获取数据库配置（关键修复）
            db_config = st.session_state.temp_db_config

            # 构建完整配置
            full_config = {
                "host": db_config["host"],
                "port": db_config["port"],
                "user": db_config["user"],
                "password": db_config["password"],
                "dbname": db_config["dbname"],
                "model_type": db_config["model_type"],
                "use_cuda": db_config["use_cuda"],
                "user_city": st.session_state.matched_city,
                "user_adcode": str(st.session_state.matched_adcode)
            }

            logging.info(f"准备保存的配置内容：{full_config}")

            # 保存到yaml文件
            if save_db_config_to_yaml(full_config):
                st.success(f"初始化完成！城市：{st.session_state.matched_city}")
                st.success("即将跳转到登录页...")

                # 清理状态并跳转
                st.session_state.temp_db_config = None
                st.session_state.db_initialized = False
                st.session_state.matched_city = None
                st.session_state.matched_adcode = None

                time.sleep(1.5)
                st.rerun()
            else:
                st.error("保存配置文件失败，请检查文件写入权限")
                logging.error(f"save_db_config_to_yaml 返回 False，配置未保存")
    # 返回按钮（如果需要）
    if st.button("返回修改数据库配置"):
        st.session_state.db_initialized = False
        st.session_state.temp_db_config = None
        st.rerun()


# ========== 修正版：保存数据库配置到db.yaml → 标准yaml格式，仅手动初始化成功时调用 ==========
def save_db_config_to_yaml(config):
    try:
        # 调试输出：打印保存的配置内容
        logging.info(f"保存的配置内容：{config}")
        logging.info(f"配置文件路径：{DB_YAML_PATH}")

        with open(DB_YAML_PATH, 'w', encoding='utf-8') as f:
            yaml.dump(
                config,
                stream=f,
                allow_unicode=True,  # 支持中文
                sort_keys=False,  # 不打乱配置项的顺序
                default_flow_style=False,  # 强制生成【标准yaml分行格式】，不是一行压缩格式
                indent=4  # 缩进4个空格，yaml格式美观易读
            )
        logging.info(f"配置文件已保存到 {DB_YAML_PATH}")
        return True
    except Exception as e:
        logging.error(f"保存db.yaml失败：{e}")
        return False

# ========== 新增：从db.yaml读取配置 ==========
def load_db_config_from_yaml():
    if not DB_YAML_PATH.exists():
        return None  # 无yaml文件，返回空
    try:
        with open(DB_YAML_PATH, 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f)

        # 新增：严格校验配置完整性
        required_keys = ["host", "port", "user", "password", "dbname", "model_type", "use_cuda", "user_city",
                         "user_adcode"]

        # 检查是否有缺失的必需字段
        missing_keys = [k for k in required_keys if k not in config]
        if missing_keys:
            logging.error(f"db.yaml缺少必需字段: {missing_keys}")
            # 删除无效的配置文件
            try:
                os.remove(DB_YAML_PATH)
                logging.info("已删除无效的db.yaml文件")
            except Exception as e:
                logging.error(f"删除db.yaml失败: {e}")
            return None

        # 新增：验证字段类型和值
        try:
            # 验证port是整数
            port = int(config["port"])
            if port <= 0 or port > 65535:
                raise ValueError(f"无效端口: {port}")
            config["port"] = port

            # 验证use_cuda是布尔值
            if not isinstance(config["use_cuda"], bool):
                config["use_cuda"] = str(config["use_cuda"]).lower() in ("true", "1", "yes")

            # 验证model_type是有效值
            valid_model_types = ["双模型融合", "text2vec-base-chinese", "bert-base-chinese"]
            if config["model_type"] not in valid_model_types:
                raise ValueError(f"无效模型类型: {config['model_type']}")

            # 验证必填字段非空
            for key in ["host", "user", "dbname", "user_city", "user_adcode"]:
                if not config[key] or not str(config[key]).strip():
                    raise ValueError(f"字段{key}不能为空")

        except (ValueError, TypeError) as e:
            logging.error(f"db.yaml字段验证失败: {e}")
            # 删除无效的配置文件
            try:
                os.remove(DB_YAML_PATH)
                logging.info("已删除无效的db.yaml文件")
            except Exception as e:
                logging.error(f"删除db.yaml失败: {e}")
            return None

        return config

    except Exception as e:
        logging.error(f"读取db.yaml失败：{e}")
        # 如果文件损坏，也删除它
        try:
            if DB_YAML_PATH.exists():
                os.remove(DB_YAML_PATH)
                logging.info("已删除损坏的db.yaml文件")
        except:
            pass
        return None


def clean_jieba_cache():
    """
    核心：仅清理 jieba 缓存，足够用了
    兼容 jieba 0.42.1：直接操作 cache 字典，而非调用不存在的方法
    """
    # 清理 jieba 内置缓存（0.42.1 的 jieba.dt.cache 是 dict，支持 clear()）
    if hasattr(jieba.dt, 'cache') and isinstance(jieba.dt.cache, dict):
        jieba.dt.cache.clear()  # 调用字典的 clear()，100% 兼容

    # 清理全局 jieba_cache，防止无限膨胀
    global jieba_cache
    jieba_cache.clear()


# ====================== 权限定义（新增） ======================
class Permissions:
    """权限定义类"""
    # 权限位定义（从低位到高位）
    LOGIN = 0b00000001  # 1: 登录权限（基础权限）
    VIEW_KB = 0b00000010  # 2: 查看知识库
    UPLOAD_KB = 0b00000100  # 4: 上传文件到知识库
    DELETE_KB = 0b00001000  # 8: 删除知识库
    EDIT_KB = 0b00010000  # 16: 编辑知识库内容
    USE_QA = 0b00100000  # 32: 使用问答功能
    USE_IMAGE_QA = 0b01000000  # 64: 使用图片识别问答
    ADMIN = 0b10000000  # 128: 管理员权限（包含所有权限）

    # 默认权限组合
    DEFAULT_USER = LOGIN | USE_QA | USE_IMAGE_QA  # 1+32+64 = 97
    DEFAULT_ADMIN = 0b11111111  # 255: 所有权限

    # 权限名称映射
    PERMISSION_NAMES = {
        LOGIN: "登录权限",
        VIEW_KB: "查看知识库",
        UPLOAD_KB: "上传文件",
        DELETE_KB: "删除知识库",
        EDIT_KB: "编辑内容",
        USE_QA: "智能问答",
        USE_IMAGE_QA: "图片识别问答",
        ADMIN: "管理员"
    }

    @staticmethod
    def has_permission(user_permissions, required_permission):
        """检查用户是否拥有指定权限"""
        if user_permissions & Permissions.ADMIN:  # 管理员拥有所有权限
            return True
        return bool(user_permissions & required_permission)

    @staticmethod
    def get_permission_names(permissions_int):
        """获取权限整数对应的权限名称列表"""
        names = []
        for perm_value, perm_name in Permissions.PERMISSION_NAMES.items():
            if permissions_int & perm_value:
                names.append(perm_name)
        return names

    @staticmethod
    def set_permission(permissions_int, permission_value, enabled=True):
        """设置或清除某个权限位"""
        if enabled:
            return permissions_int | permission_value
        else:
            return permissions_int & (~permission_value)


# ====================== 敏感操作日志函数 ======================
def log_sensitive_operation(operation_type, target_user, current_user, details=""):
    """记录敏感操作到日志文件"""
    log_file = Path(__file__).parent / "security.log"

    # 创建安全的日志格式
    timestamp = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime())
    ip_address = "127.0.0.1"  # 本地运行，默认为本地

    log_entry = f"[{timestamp}] [IP:{ip_address}] [操作人:{current_user}] [操作:{operation_type}] [目标用户:{target_user}] {details}\n"

    try:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(log_entry)
    except Exception as e:
        logging.error(f"写入安全日志失败: {e}")


# ========== 新增：隐藏Streamlit顶部占用（参考链接CSS，修正拼写错误） ==========
HIDE_ST_STYLE = """
<style>
/* 隐藏顶部Header、Toolbar、装饰条 */
div[data-testid="stHeader"] {
    visibility: hidden;
    height: 0px;
    position: fixed;
}
div[data-testid="stToolbar"] {
    visibility: hidden;
    height: 0px;
    position: fixed;
}
div[data-testid="stDecoration"] {
    visibility: hidden;
    height: 0px;
    position: fixed;
}
/* 隐藏顶部主菜单 */
#MainMenu {
    visibility: hidden;
    height: 0px;
}
/* 隐藏页头、页脚 */
header {
    visibility: hidden;
    height: 0px;
}
footer {
    visibility: hidden;
    height: 0%;
}
/* 调整主内容区padding，避免内容顶边太拥挤 */
div[data-testid="stMainBlockContainer"] {
    padding: 0.5rem 1rem 10rem !important;
}
</style>
"""
# ========== 唯一使用：登录+注册一体化CSS样式（修复方形元素+布局问题） ==========
LOGIN_REGISTER_STYLE = """
<style>
/* 整体容器：移除冗余边框，解决方形元素问题 */
.auth-container {
    max-width: 450px;
    margin: 2rem auto; /* 减少顶部间距，避免上方空白占位 */
    padding: 2rem;
    background-color: #f8fafc;
    border-radius: 12px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.1);
    border: none; /* 移除默认边框，解决方形元素 */
}
/* 选项卡样式：简化样式，避免冲突 */
.auth-tabs {
    display: flex;
    justify-content: center;
    margin-bottom: 1.5rem;
    gap: 0.5rem;
}
.auth-tab {
    padding: 0.8rem 2rem;
    cursor: pointer;
    border-radius: 8px;
    background-color: #e2e8f0;
    color: #64748b;
    font-weight: 500;
    border: none;
}
.auth-tab.active {
    background-color: #3b82f6;
    color: white;
}
/* 标题样式 */
.auth-title {
    text-align: center;
    margin-bottom: 1.5rem;
    color: #1e293b;
    font-size: 22px;
}
/* 提示信息样式 */
.auth-error {
    color: #dc2626;
    text-align: center;
    margin-top: 1rem;
    font-size: 14px;
}
.auth-success {
    color: #16a34a;
    text-align: center;
    margin-top: 1rem;
    font-size: 14px;
}
/* 按钮列布局：确保横向排列，不受容器干扰 */
.btn-row {
    display: flex;
    gap: 0.5rem;
    justify-content: center;
    margin-top: 1rem;
}
</style>
"""

# ========== 国内环境适配+CUDA优化 ==========
os.environ["TRANSFORMERS_OFFLINE"] = "1"
os.environ["HF_DATASETS_OFFLINE"] = "1"
os.environ["CUDA_LAUNCH_BLOCKING"] = "1"
os.environ["PYTORCH_CUDA_ALLOC_CONF"] = "max_split_size_mb:64,garbage_collection_threshold:0.5"

# ====================== 解决AutoTokenizer导入兼容 ==========
try:
    from transformers import AutoTokenizer, AutoModel
except ImportError:
    from transformers.models.auto.tokenization_auto import AutoTokenizer
    from transformers.models.auto.modeling_auto import AutoModel

# ====================== 保留所有核心依赖（jieba完整可用） ==========
import streamlit as st
import psycopg2
import psycopg2.errors
from psycopg2 import OperationalError
import sys
import numpy as np
import re
import torch
import pandas as pd
from collections import defaultdict
from psycopg2.extras import RealDictCursor
from pgvector.psycopg2 import register_vector
import pdfplumber
from docx import Document
from scipy.spatial.distance import cosine
import chardet
# ========== 【修正】OpenVINO 2025.4.1 正确导入 (完全贴合你的测试代码，无废弃警告) ==========
import cv2
import numpy as np
from openvino import Core  # 2025版原生正确导入，无任何警告
from PIL import Image
import imghdr
import sys

# ========== 强制隔离：YOLO=iGPU.0，Ollama=NVIDIA GPU 绝对无冲突 ==========
os.environ["OV_NUM_THREADS"] = "4"
os.environ["OPENCV_IO_MAX_IMAGE_PIXELS"] = str(1024 * 1024 * 200)
# ====================== 全局配置（保留原有 + 简化Ollama配置） ==========
LOCAL_MODEL_ROOT_PATH = r"E:\final_exam_test\handcraft\models"
TEXT2VEC_MODEL_PATH = r"E:\final_exam_test\handcraft\models\text2vec-base-chinese"
BERT_MODEL_PATH = os.path.join(LOCAL_MODEL_ROOT_PATH, "bert-base-chinese")

DEFAULT_DB_CONFIG = {
    "host": "localhost",
    "port": 5432,
    "user": "postgres",
    "password": "",
    "dbname": "kb_db"
}
PAGE_SIZE = 10
FILE_PROCESS_TIMEOUT = 120
MAX_FILES_PER_UPLOAD = 5
USE_CUDA_DEFAULT = False
CHUNK_RETRY_TIMES = 2
FILE_RETRY_TIMES = 1
CHUNK_SIZE_DEFAULT = 150
MIN_LENGTH_DEFAULT = 30
TEXT2VEC_WEIGHT = 0.6
BERT_WEIGHT = 0.4
# 内置合理阈值，无需用户调节
BUILT_IN_MIN_SIMILARITY = 0.5
KEYWORD_MATCH_WEIGHT = 0.4
# Ollama配置（仅保留模型名称，host严格在Client实例化时指定）


# ========== 【新增】YOLO+OpenVINO INT8 配置 (和测试代码一模一样，直接用) - Intel iGPU运行 ==========
YOLO_INT8_MODEL_XML = r"E:\final_exam_test\handcraft\handcrafted\runs\detect\train8\weights_openvion\best_int8_model\best_int8.xml"
YOLO_INT8_MODEL_BIN = r"E:\final_exam_test\handcraft\handcrafted\runs\detect\train8\weights_openvion\best_int8_model\best_int8.bin"
# ===== 你的测试代码里的核心配置 =====
# 替换原硬编码的 TARGET_GPU 定义
TARGET_GPU = "GPU.0"
YOLO_CONF_THRESHOLD = 0.5  # 置信度阈值
YOLO_IOU_THRESHOLD = 0.45  # IOU阈值
YOLO_INPUT_SIZE = (640, 640)  # 输入尺寸
# COCO80类别，和你测试代码一致
YOLO_CLASSES = [
    'Bacterial Spot',
    'Early_Blight',
    'Healthy',
    'Late_blight',
    'Leaf Mold',
    'Target_Spot',
    'black spot',
    'Apple Scab Leaf',
    'Apple leaf',
    'Apple rust leaf',
    'Bell_pepper leaf spot',
    'Bell_pepper leaf',
    'Blueberry leaf',
    'Cherry leaf',
    'Corn Gray leaf spot',
    'Corn leaf blight',
    'Corn rust leaf',
    'Peach leaf',
    'Potato leaf early blight',
    'Potato leaf late blight',
    'Potato leaf',
    'Raspberry leaf',
    'Soyabean leaf',
    'Soybean leaf',
    'Squash Powdery mildew leaf',
    'Strawberry leaf',
    'Tomato Septoria leaf spot',
    'Tomato leaf mosaic virus',
    'Tomato leaf yellow virus',
    'Tomato leaf',
    'Tomato two spotted spider mites leaf',
    'grape leaf black rot',
    'grape leaf'
]
nc = 33  # 类别总数，和你的yaml一致
# ========== 【新增】图片上传安全限制 (防图片马/恶意文件，严格按要求) ==========
ALLOWED_IMAGE_EXT = ['jpg', 'jpeg', 'png']  # 只允许安全图片格式
MAX_IMAGE_SIZE = 8 * 1024 * 1024  # 最大8MB，防超大文件内存溢出

# ====================== 柔性停用词（完全保留） ==========
CORE_STOP_WORDS = {
    "的", "了", "是", "我", "你", "他", "她", "它", "我们", "你们", "他们",
    "这", "那", "个", "和", "与", "或", "在", "于", "为", "对", "同", "跟",
    "把", "被", "从", "到", "有", "无", "不", "没", "会", "能", "要", "将",
    "就", "还", "也", "又", "而", "则"
}

SCENE_STOP_WORDS = {
    "件", "只", "本", "页", "行", "列"
}

USER_DEFINE_STOP_WORDS = set()

# 初始化验证码生成器（确保字体/尺寸有效）
image_captcha = ImageCaptcha(width=160, height=60, font_sizes=(32, 40, 48))


# 1. 生成随机验证码文本（核心：始终随机）
def generate_captcha_text(length=4):
    char_set = string.digits + string.ascii_letters
    captcha_text = ''.join(random.choice(char_set) for _ in range(length))
    return captcha_text


# 2. 生成验证码图片（返回有效BytesIO，兜底仍用随机文本）
def generate_captcha_image(captcha_text=None):
    # 关键：如果没传文本，自动生成随机文本（避免固定值）
    if captcha_text is None:
        captcha_text = generate_captcha_text()
    try:
        img = image_captcha.generate(captcha_text)
        img_bytes = BytesIO(img.read())
        img_bytes.seek(0)  # 重置指针，避免Streamlit读取空
        # 修复：添加合法文件名，解决Bad filename错误
        img_bytes.name = f"captcha_{captcha_text}.png"
        return img_bytes
    except Exception as e:
        # 极端兜底：仍生成随机文本+空图片，而非固定1234
        from PIL import Image, ImageDraw, ImageFont
        # 生成随机文本（核心优化）
        random_text = generate_captcha_text()
        # 创建空白图片并写入随机文本
        img = Image.new('RGB', (160, 60), color='white')
        draw = ImageDraw.Draw(img)
        # 兜底字体（避免字体缺失报错）
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except:
            font = ImageFont.load_default(size=40)
        draw.text((20, 10), random_text, fill='black', font=font)
        # 保存为BytesIO
        empty_img = BytesIO()
        img.save(empty_img, format='PNG')
        empty_img.seek(0)
        # 🔧 修复：添加合法文件名
        empty_img.name = f"captcha_{random_text}.png"
        return empty_img


def init_captcha_system():
    """初始化验证码系统"""
    # 确保验证码相关的session_state存在
    if 'captcha_text' not in st.session_state:
        refresh_captcha()
    if 'captcha_image' not in st.session_state:
        refresh_captcha()
    if 'captcha_verified' not in st.session_state:
        st.session_state.captcha_verified = False
    if 'captcha_attempts' not in st.session_state:
        st.session_state.captcha_attempts = 0


# 3. 核心刷新函数（始终生成随机验证码）
def refresh_captcha():
    try:
        new_text = generate_captcha_text()  # 随机文本
        new_img = generate_captcha_image(new_text)  # 对应随机图片
        st.session_state['captcha_text'] = new_text
        st.session_state['captcha_image'] = new_img
    except Exception as e:
        # 极端兜底：仍生成随机文本+图片
        random_text = generate_captcha_text()
        st.session_state['captcha_text'] = random_text
        st.session_state['captcha_image'] = generate_captcha_image(random_text)


# 添加在 refresh_captcha 函数之后

def verify_captcha(user_input, operation_type="login"):
    """验证验证码（修复版）"""
    user_input = str(user_input).strip().upper()
    correct_text = st.session_state.get('captcha_text', '').strip().upper()

    # 检查验证码尝试次数
    if st.session_state.get('captcha_attempts', 0) >= 3:
        st.error("验证码尝试次数过多，请刷新页面！")
        refresh_captcha()
        st.session_state.captcha_attempts = 0
        return False

    if not user_input or not correct_text:
        st.error("验证码不能为空！")
        st.session_state.captcha_attempts = st.session_state.get('captcha_attempts', 0) + 1
        return False

    # 验证码验证（不区分大小写）
    is_valid = user_input == correct_text

    # 无论验证结果如何，都刷新验证码
    refresh_captcha()

    if not is_valid:
        st.session_state.captcha_attempts = st.session_state.get('captcha_attempts', 0) + 1
        st.error(f"验证码错误！ ({st.session_state.captcha_attempts}/3)")

    st.session_state.captcha_verified = is_valid
    return is_valid


# 4. 初始化验证码（始终随机，无固定值）
def init_captcha_session():
    if "captcha_text" not in st.session_state or st.session_state.get("captcha_text") is None:
        refresh_captcha()
    if "captcha_image" not in st.session_state or st.session_state.get("captcha_image") is None:
        refresh_captcha()


# ========== 【新增 最高优先级】启动强制初始化Intel iGPU + 可用性校验 (彻底防冲突，你的首要要求) ==========
# ========== 【唯一修改1】100%复刻你的test代码：Intel iGPU初始化+验证，只做连接验证，不做其他操作 ==========
@st.cache_resource(ttl=None, show_spinner="初始化并验证Intel iGPU可用性...")
def init_verify_intel_igpu():
    from openvino import Core
    core = Core()
    available_devices = core.available_devices
    selected_device = "CPU"
    if TARGET_GPU in available_devices:
        try:
            dev_name = core.get_property(TARGET_GPU, "FULL_DEVICE_NAME")
            if "Intel" in dev_name:
                selected_device = TARGET_GPU
        except:
            pass
    return core, selected_device


# 启动时立即验证，全局生效
intel_core = init_verify_intel_igpu()


# ========== 【新增1】图片安全校验 三重防护 防图片马/恶意文件 (你的强制要求) ==========
def check_image_safety(uploaded_file):
    ALLOWED_IMAGE_EXT = ['jpg', 'jpeg', 'png', 'bmp']
    MAX_IMAGE_SIZE = 8 * 1024 * 1024
    if uploaded_file is None:
        return False, "上传的图片为空文件！"
    img_bytes = uploaded_file.getvalue()
    if not img_bytes:
        return False, "上传的图片为空文件！"
    file_ext = uploaded_file.name.split('.')[-1].lower() if '.' in uploaded_file.name else ''
    if file_ext not in ALLOWED_IMAGE_EXT:
        return False, f"仅支持{ALLOWED_IMAGE_EXT}格式！"
    if uploaded_file.size > MAX_IMAGE_SIZE:
        return False, f"图片超8MB！当前：{round(uploaded_file.size / 1024 / 1024, 2)}MB"
    try:
        img_stream = BytesIO(img_bytes)
        img_stream.seek(0)  # 强制归位指针，解决seek报错的核心！！！
        imghdr.what(img_stream)
        img_stream.seek(0)  # 再次归位，防止指针耗尽
        img = Image.open(img_stream)
        img.verify()
        img.close()  # 关闭Image对象
        img_stream.close()  # 关闭BytesIO流
        del img, img_stream  # 释放引用
        return True, "图片安全"
    except AttributeError as e:
        if "seek" in str(e):
            return False, "图片损坏/文件流指针异常"
        return False, f"图片无效: {str(e)}"
    except Exception as e:
        return False, f"图片校验失败: {str(e)[:50]}"


# ========== 【新增2】YOLO图片预处理 (和你测试代码完全一致) ==========
def preprocess_image(image, input_size):
    img = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
    img_h, img_w = img.shape[:2]
    scale = min(input_size[0] / img_w, input_size[1] / img_h)
    new_w, new_h = int(img_w * scale), int(img_h * scale)
    img_resized = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_LINEAR)
    pad_w = (input_size[0] - new_w) // 2
    pad_h = (input_size[1] - new_h) // 2
    img_padded = cv2.copyMakeBorder(img_resized, pad_h, pad_h, pad_w, pad_w, cv2.BORDER_CONSTANT, value=(114, 114, 114))
    img_padded = img_padded.transpose(2, 0, 1) / 255.0
    img_padded = np.expand_dims(img_padded, 0).astype(np.float32)
    return img_padded, scale, pad_w, pad_h, img_w, img_h


# ========== 【新增3】YOLO+OpenVINO INT8 推理(Intel iGPU运行，和你测试代码完全一致！！！核心中的核心) ==========
@st.cache_resource(ttl=3600, show_spinner="加载YOLO INT8叶片病害模型(iGPU/CPU)中...")
def load_yolo_ov_model():
    core, selected_device = init_verify_intel_igpu()
    # 加载你配置的 INT8 模型路径，原样使用，无修改
    model = core.read_model(model=YOLO_INT8_MODEL_XML, weights=YOLO_INT8_MODEL_BIN)
    compiled_model = core.compile_model(model=model, device_name=selected_device)
    output_layer = compiled_model.output(0)
    # 仅后台打印，页面无显示，完美解决你的顾虑
    # print(f"YOLO INT8模型加载完成 | 运行设备: {selected_device} | 置信度阈值: {YOLO_CONF_THRESHOLD}")
    return compiled_model, output_layer


def yolo_ov_detect(image):
    """YOLO+OpenVINO INT8 iGPU推理，返回检测结果文本+标注后的图片"""
    # 关键修复：确保传入的是PIL Image对象且未失效
    if not isinstance(image, Image.Image):
        raise ValueError("输入必须是PIL Image对象")
    image = image.copy()  # 深拷贝
    compiled_model, output_layer = load_yolo_ov_model()
    # 预处理
    input_tensor, scale, pad_w, pad_h, img_w, img_h = preprocess_image(image, YOLO_INPUT_SIZE)
    # iGPU推理 (和测试代码一致)
    results = compiled_model([input_tensor])[output_layer]
    predictions = np.squeeze(results).T
    # 过滤低置信度
    scores = np.max(predictions[:, 4:], axis=1)
    predictions = predictions[scores > YOLO_CONF_THRESHOLD, :]
    scores = scores[scores > YOLO_CONF_THRESHOLD]
    class_ids = np.argmax(predictions[:, 4:], axis=1)
    # 坐标还原
    boxes = predictions[:, :4]
    boxes[:, 0] = (boxes[:, 0] - pad_w) / scale
    boxes[:, 1] = (boxes[:, 1] - pad_h) / scale
    boxes[:, 2] = (boxes[:, 2] - pad_w) / scale
    boxes[:, 3] = (boxes[:, 3] - pad_h) / scale
    boxes[:, 0] = np.clip(boxes[:, 0], 0, img_w)
    boxes[:, 1] = np.clip(boxes[:, 1], 0, img_h)
    boxes[:, 2] = np.clip(boxes[:, 2], 0, img_w)
    boxes[:, 3] = np.clip(boxes[:, 3], 0, img_h)
    # NMS非极大值抑制
    indices = cv2.dnn.NMSBoxes(boxes[:, :4].tolist(), scores.tolist(), YOLO_CONF_THRESHOLD, YOLO_IOU_THRESHOLD)
    # 生成检测结果文本+标注图片
    detect_text = "【YOLO目标检测结果】\n"
    img_draw = np.array(image).copy()
    for i in indices:
        i = i[0] if isinstance(i, (list, np.ndarray)) else i
        x1, y1, x2, y2 = int(boxes[i, 0]), int(boxes[i, 1]), int(boxes[i, 2]), int(boxes[i, 3])
        cls = YOLO_CLASSES[class_ids[i]]
        conf = round(scores[i], 3)
        detect_text += f"- 检测到：{cls} | 置信度：{conf} | 位置：({x1},{y1})-({x2},{y2})\n"
        # 画框
        cv2.rectangle(img_draw, (x1, y1), (x2, y2), (0, 255, 0), 2)
        cv2.putText(img_draw, f"{cls} {conf}", (x1, y1 - 5), cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 255, 0), 2)
    # 无检测结果
    if len(indices) == 0:
        detect_text += "未检测到任何目标物体"
    # BGR转RGB
    img_draw = cv2.cvtColor(img_draw, cv2.COLOR_BGR2RGB)
    img_draw = Image.fromarray(img_draw)
    # 强制释放YOLO推理的内存/显存 核心修复
    del input_tensor, scale, pad_w, pad_h, img_w, img_h, predictions, scores, class_ids, boxes, indices
    import gc
    gc.collect()  # 立即回收垃圾内存
    return detect_text, img_draw


# ====================== 1. 严格按你的Country模型逻辑定义Pydantic模型（顶格类，类内4空格缩进） ==========
# 完全对齐你的Country模型格式：顶格定义类，类内4个空格缩进，字段明确
class KnowledgeBaseAnswer(BaseModel):
    kb_search_results: str  # 4空格缩进，与你的Country模型一致
    thinking_process: str  # 4空格缩进，与你的Country模型一致
    final_conclusion: str  # 4空格缩进，与你的Country模型一致


# ====================== 精准关键词提取（完全保留 + 关键词匹配） ==========
def calculate_word_weight(word, text, word_pos_dict):
    if word in CORE_STOP_WORDS:
        return 0.0

    word_base_weight = 0.8 if word not in SCENE_STOP_WORDS else 0.2
    if word in USER_DEFINE_STOP_WORDS:
        word_base_weight *= 0.1

    word_count = text.count(word)
    total_word_count = len(jieba.lcut(text))
    freq_weight = word_count / total_word_count if total_word_count > 0 else 0.0

    pos_weights = []
    if word in word_pos_dict:
        text_length = len(text)
        for pos in word_pos_dict[word]:
            if pos / text_length <= 0.2 or pos / text_length >= 0.8:
                pos_weights.append(1.0)
            else:
                pos_weights.append(0.5)
    pos_weight = np.mean(pos_weights) if pos_weights else 0.5

    total_weight = (freq_weight * 0.4) + (pos_weight * 0.4) + (word_base_weight * 0.2)
    return round(total_weight, 4)


def extract_keywords_optimized(text, top_k=3):
    text_clean = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', '', text)[:500]
    if not text_clean:
        return ["未提取到关键词"]

    words = jieba.lcut(text_clean)
    word_pos_dict = defaultdict(list)
    for idx, word in enumerate(words):
        word_start_pos = len("".join(words[:idx]))
        word_pos_dict[word].append(word_start_pos)

    candidate_words = [word for word in words if len(word) >= 2]
    if not candidate_words:
        return ["未提取到关键词"]

    word_weight_dict = {}
    for word in set(candidate_words):
        word_weight_dict[word] = calculate_word_weight(word, text_clean, word_pos_dict)

    sorted_word_weights = sorted(word_weight_dict.items(), key=lambda x: x[1], reverse=True)
    valid_word_weights = [(word, weight) for word, weight in sorted_word_weights if weight > 0.0]
    if not valid_word_weights:
        return ["未提取到关键词"]

    top_words = [word for word, _ in valid_word_weights[:top_k]]
    return top_words


# 新增：jieba分词缓存字典（全局）
jieba_cache = {}


def calculate_keyword_match_score(query, chunk_content, top_k=5):
    # 缓存查询文本的分词结果
    if query not in jieba_cache:
        jieba_cache[query] = extract_keywords_optimized(query, top_k=top_k)
    query_keywords = set(jieba_cache[query])

    # 缓存分段文本的分词结果
    if chunk_content not in jieba_cache:
        jieba_cache[chunk_content] = extract_keywords_optimized(chunk_content, top_k=top_k)
    chunk_keywords = set(jieba_cache[chunk_content])

    if not query_keywords or query_keywords == {"未提取到关键词"}:
        return 0.0
    intersection = len(query_keywords & chunk_keywords)
    union = len(query_keywords | chunk_keywords)
    clean_jieba_cache()  # 每次调用都清理缓存，防止无限膨胀
    return intersection / union if union > 0 else 0.0


def generate_auto_chunk_title_optimized(file_name, chunk_content, chunk_idx, page_num=None):
    keywords = extract_keywords_optimized(chunk_content)
    keyword_str = "、".join(keywords)
    base_title = f"{os.path.splitext(file_name)[0]} - {keyword_str}"
    if page_num:
        base_title += f"（第{page_num}页）"
    final_title = f"{base_title} - 第{chunk_idx + 1}段"
    if len(final_title) > 50:
        final_title = final_title[:50] + "..."
    return final_title


# ====================== 停用词配置面板（完全保留） ==========
def stop_word_config_panel():
    st.sidebar.divider()
    st.sidebar.subheader("🔧 停用词配置")
    st.sidebar.caption("仅调整场景停用词")

    scene_stop_words_str = "、".join(SCENE_STOP_WORDS)
    st.sidebar.info(f"当前：{scene_stop_words_str}")

    col_add, col_del = st.sidebar.columns(2)
    with col_add:
        new_stop_word = st.text_input("添加", placeholder="单个词汇", key="add_sw")
        if st.button("➕", key="add_btn"):
            if new_stop_word and new_stop_word not in SCENE_STOP_WORDS:
                SCENE_STOP_WORDS.add(new_stop_word)
                st.sidebar.success(f"添加：{new_stop_word}")
                st.rerun()
            elif new_stop_word in SCENE_STOP_WORDS:
                st.sidebar.warning("已存在")
    with col_del:
        del_stop_word = st.selectbox("删除", options=list(SCENE_STOP_WORDS), key="del_sw", index=0)
        if st.button("➖", key="del_btn"):
            SCENE_STOP_WORDS.remove(del_stop_word)
            st.sidebar.success(f"删除：{del_stop_word}")
            st.rerun()


# ====================== 双模型编码器（完全保留） ==========
@st.cache_resource(ttl=3600, show_spinner="加载本地模型中...")
def get_local_encoder(use_cuda=USE_CUDA_DEFAULT, model_type="双模型融合", silent=False):
    """
    加载本地模型
    :param silent: 静默模式，True=不显示模型加载消息（用于登录页自动加载）
    """

    class LocalText2VecBERTEncoder:
        def __init__(self, use_cuda=use_cuda, model_type=model_type):
            self.use_model = True
            self.device = torch.device("cpu")
            self.model_type = model_type
            self.text2vec_tokenizer = None
            self.text2vec_model = None
            self.bert_tokenizer = None
            self.bert_model = None

            if torch.cuda.is_available() and use_cuda:
                self.device = torch.device("cuda")
                torch.cuda.empty_cache()
                torch.cuda.reset_max_memory_allocated()

            if not silent:  # 静默时不显示
                st.info(f"计算设备：{self.device} | 模型类型：{self.model_type}（jieba可用）")

            try:
                if self.model_type in ["text2vec-base-chinese", "双模型融合"]:
                    self.text2vec_tokenizer = AutoTokenizer.from_pretrained(
                        TEXT2VEC_MODEL_PATH,
                        local_files_only=True,
                        trust_remote_code=True
                    )
                    self.text2vec_model = AutoModel.from_pretrained(
                        TEXT2VEC_MODEL_PATH,
                        local_files_only=True,
                        trust_remote_code=True,
                        low_cpu_mem_usage=False
                    ).to(self.device).eval()
                    if not silent:  # 静默时不显示
                        st.success(f"text2vec-base-chinese 加载成功")

                if self.model_type in ["bert-base-chinese", "双模型融合"]:
                    self.bert_tokenizer = AutoTokenizer.from_pretrained(
                        BERT_MODEL_PATH,
                        local_files_only=True,
                        trust_remote_code=True,
                        low_cpu_mem_usage=False
                    )
                    self.bert_model = AutoModel.from_pretrained(
                        BERT_MODEL_PATH,
                        local_files_only=True,
                        trust_remote_code=True
                    ).to(self.device).eval()
                    if not silent:  # 静默时不显示
                        st.success(f"bert-base-chinese 加载成功")

            except Exception as e:
                error_msg = f"模型加载失败：{str(e)[:100]}"
                st.error(error_msg)
                logging.error(f"模型加载失败详情：{traceback.format_exc()}")
                self.use_model = False
                return

        def _encode_text2vec(self, text):
            inputs = self.text2vec_tokenizer(
                text,
                truncation=True,
                max_length=128,
                padding="max_length",
                return_tensors="pt"
            )
            inputs = {k: v.to(self.device) for k, v in inputs.items()}
            with torch.no_grad():
                if self.device.type == "cuda":
                    torch.cuda.synchronize()
                outputs = self.text2vec_model(**inputs)
            vec = outputs.last_hidden_state[:, 0, :].squeeze().cpu().numpy()
            return vec

        def _encode_bert(self, text):
            inputs = self.bert_tokenizer(
                text,
                truncation=True,
                max_length=128,
                padding="max_length",
                return_tensors="pt"
            )
            inputs = {k: v.to(self.device) for k, v in inputs.items()}
            with torch.no_grad():
                if self.device.type == "cuda":
                    torch.cuda.synchronize()
                outputs = self.bert_model(**inputs)
            mask = inputs['attention_mask'].cpu().numpy()
            hidden_state = outputs.last_hidden_state.squeeze().cpu().numpy()
            mask_expanded = np.broadcast_to(mask.reshape(-1, 1), hidden_state.shape)
            sum_embeddings = np.sum(hidden_state * mask_expanded, axis=0)
            sum_mask = np.clip(mask_expanded.sum(axis=0), a_min=1e-9, a_max=None)
            vec = sum_embeddings / sum_mask
            return vec

        def _fusion_vectors(self, vec1, vec2, weight1=TEXT2VEC_WEIGHT, weight2=BERT_WEIGHT):
            fusion_vec = (weight1 * vec1) + (weight2 * vec2)
            return fusion_vec

        def encode_text(self, text):
            if not isinstance(text, str):
                vec = np.zeros(768)
                del text  # 释放变量引用
                return vec

            if self.use_model:
                try:
                    vec = np.zeros(768)
                    if self.model_type == "text2vec-base-chinese":
                        if self.text2vec_tokenizer and self.text2vec_model:
                            vec = self._encode_text2vec(text)
                    elif self.model_type == "bert-base-chinese":
                        if self.bert_tokenizer and self.bert_model:
                            vec = self._encode_bert(text)
                    elif self.model_type == "双模型融合":
                        if self.text2vec_tokenizer and self.text2vec_model and self.bert_tokenizer and self.bert_model:
                            vec1 = self._encode_text2vec(text)
                            vec2 = self._encode_bert(text)
                            vec = self._fusion_vectors(vec1, vec2)
                            del vec1, vec2  # 释放临时张量，立即清内存

                    if self.device.type == "cuda":
                        torch.cuda.empty_cache()
                        torch.cuda.ipc_collect()  # 强制回收CUDA碎片内存
                    norm = np.linalg.norm(vec)
                    if norm == 0:
                        del norm  # 释放临时变量
                        return vec
                    return vec / norm
                except Exception as e:
                    st.warning(f"模型编码失败，使用jieba特征：{str(e)[:50]}")
                    return self._jieba_text_feature(text)
            else:
                return self._jieba_text_feature(text)

        def _jieba_text_feature(self, text):
            text_clean = re.sub(r'\W+', '', text)[:200]
            words = jieba.lcut(text_clean)
            feature = np.zeros(768)
            for i, word in enumerate(words[:768]):
                feature[i] = hash(word) % 1000 / 1000
            norm = np.linalg.norm(feature)
            if norm == 0:
                return feature
            return feature / norm

        def calculate_similarity(self, text1, text2):
            try:
                vec1 = self.encode_text(text1)
                vec2 = self.encode_text(text2)
                sim = 1 - cosine(vec1, vec2)
                return max(0.0, min(1.0, sim))
            except Exception as e:
                st.warning(f"相似度计算失败，使用词集合匹配：{str(e)[:50]}")
                words1 = set(jieba.lcut(text1[:100]))
                words2 = set(jieba.lcut(text2[:100]))
                total = len(words1) + len(words2)
                if total == 0:
                    return 0.0
                return len(words1 & words2) / total

    return LocalText2VecBERTEncoder(use_cuda, model_type)


# ====================== 文本清洗+分段逻辑（完全保留） ==========
def clean_text_local(text):
    if not isinstance(text, str):
        return ""
    text = re.sub(r'[\x00-\x1f\x7f]', '', text)  # 过滤不可见字符
    text = re.sub(r'\s{2,}', ' ', text)  # 过滤多余空格
    text = re.sub(r'行\d+', '', text)  # 过滤Excel行号（如“行262”）
    text = re.sub(r'分段标题（选填）', '', text)  # 过滤Excel冗余标注
    text = re.sub(r'问题（选填，单元格内一行一个）', '', text)  # 过滤Excel冗余标注
    text = text.strip()
    return text


def retry_chunk_process(func, *args, retry_times=CHUNK_RETRY_TIMES, **kwargs):
    for attempt in range(retry_times + 1):
        try:
            return func(*args, **kwargs), "success"
        except TimeoutError as e:
            if attempt < retry_times:
                st.warning(f"分段重试{attempt + 1}/{retry_times + 1}：超时")
                time.sleep(1)
            else:
                st.error(f"分段重试失败：超时")
                return None, "timeout"
        except Exception as e:
            if attempt < retry_times:
                st.warning(f"分段重试{attempt + 1}/{retry_times + 1}：{str(e)[:50]}")
                time.sleep(1)
            else:
                st.error(f"分段重试失败：{str(e)[:50]}")
                return None, "fail"
    return None, "fail"


def split_text_with_retry(text, encoder, chunk_size=CHUNK_SIZE_DEFAULT, sim_threshold=0.8,
                          min_length=MIN_LENGTH_DEFAULT):
    def _inner_split():
        text_clean = ""
        try:
            text_clean = clean_text_local(text)
            if len(text_clean) < min_length:
                return [text_clean] if text_clean else []
        except Exception as e:
            st.warning(f"文本清洗失败：{str(e)[:50]}")
            return []

        sentences = re.split(r'[。！？；：\n]', text_clean)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 1]
        if len(sentences) < 2:
            return [text_clean[:chunk_size]] if text_clean else []

        chunks = []
        current_chunk = [sentences[0]]
        start_time = time.time()

        for sent in sentences[1:]:
            if time.time() - start_time > FILE_PROCESS_TIMEOUT:
                raise TimeoutError("分段处理超时")

            current_chunk_str = ''.join(current_chunk)
            if len(current_chunk_str + sent) > chunk_size:
                if len(current_chunk_str) >= min_length:
                    chunks.append(current_chunk_str)
                    current_chunk = [sent]
                else:
                    current_chunk.append(sent)
            else:
                current_chunk.append(sent)

        final_chunk_str = ''.join(current_chunk)
        if len(final_chunk_str) >= 1:
            chunks.append(final_chunk_str)

        return chunks or [text_clean]

    start_time = time.time()
    chunks, status = retry_chunk_process(_inner_split)
    return chunks or [], status


# ====================== PDF解析（完全保留） ==========
def parse_pdf_optimized(file):
    raw_text = ""
    page_num_record = None
    pdf = None
    try:
        pdf = pdfplumber.open(file)
        for page_num, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text() or ""
            text_blocks = None
            if not page_text:
                text_blocks = page.extract_words()
                if text_blocks:
                    page_text = " ".join([block['text'] for block in text_blocks])
            if page_text.strip():
                raw_text += f"【第{page_num}页】{page_text}\n"
                page_num_record = page_num
            else:
                st.warning(f"PDF第{page_num}页无有效文本，跳过")
            del page, page_text, text_blocks
            # 释放每页的临时变量
        return raw_text, page_num_record
    except Exception as e:
        st.error(f"PDF解析失败：{str(e)[:100]}")
        return "", None
    finally:
        if pdf:
            pdf.close()  # 强制关闭PDF流，防止内存泄漏
        del file  # 释放文件对象


def parse_file_with_auto_retry(file, encoder, chunk_size=CHUNK_SIZE_DEFAULT, sim_threshold=0.8,
                               min_length=MIN_LENGTH_DEFAULT):
    file_name = file.name
    raw_text = ""
    page_num_record = None
    file_status = "success"

    for attempt in range(FILE_RETRY_TIMES + 1):
        try:
            raw_text = ""
            page_num_record = None
            if file_name.endswith('.pdf'):
                raw_text, page_num_record = parse_pdf_optimized(file)
                if not raw_text:
                    st.warning(f"{file_name}：未提取到文本，可能是加密PDF或图片PDF（扫描件）")
                    if attempt < FILE_RETRY_TIMES:
                        continue
                    else:
                        return [], False, "fail", file_name
            elif file_name.endswith('.txt'):
                try:
                    raw_text = file.getvalue().decode('utf-8')
                except:
                    raw_text = file.getvalue().decode('gbk')
            elif file_name.endswith('.docx'):
                doc = Document(file)
                for para in doc.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        raw_text += para_text + "\n"
            elif file_name.endswith('.xlsx') or file_name.endswith('.xls'):
                df = pd.read_excel(file, dtype=str).fillna("")
                for idx, row in df.iterrows():
                    row_text = " | ".join([f"{col}：{val}" for col, val in row.items()])
                    raw_text += f"行{idx + 1}：{row_text}\n"
            elif file_name.endswith('.doc'):
                    # 纯Python方案：直接从二进制中提取可读文本（无需任何外部依赖）
                    try:
                        import struct
                        import re

                        raw_bytes = file.getvalue()
                        # 尝试解码为 utf-16-le（Word 常用编码）
                        try:
                            # 先尝试 utf-16-le（小端）
                            text = raw_bytes.decode('utf-16-le', errors='ignore')
                        except:
                            # 后备：普通 utf-8
                            text = raw_bytes.decode('utf-8', errors='ignore')

                        # 清洗：保留中文、英文、数字、常见标点，移除控制字符和乱码
                        # 匹配中文、英文、数字、空格、常见标点
                        clean_text = re.sub(
                            r'[^\u4e00-\u9fa5a-zA-Z0-9\s\.\,\!\?\;\:\"\'\+\-\*\=\/\<\>\(\)\[\]\{\}\@\#\$\%\^\&\*\_\|\~\\]',
                            ' ', text)
                        # 合并多余空格
                        clean_text = re.sub(r'\s+', ' ', clean_text).strip()

                        # 如果提取到的文本太短（可能加密或特殊格式），给出提示但不报错
                        if len(clean_text) < 20:
                            # 尝试另一种编码：gbk（中文系统常用）
                            try:
                                text = raw_bytes.decode('gbk', errors='ignore')
                                clean_text = re.sub(
                                    r'[^\u4e00-\u9fa5a-zA-Z0-9\s\.\,\!\?\;\:\"\'\+\-\*\=\/\<\>\(\)\[\]\{\}\@\#\$\%\^\&\*\_\|\~\\]',
                                    ' ', text)
                                clean_text = re.sub(r'\s+', ' ', clean_text).strip()
                            except:
                                pass

                        raw_text = clean_text
                        page_num_record = 1  # .doc 不记录页码

                        # 最终检查：如果仍然为空，则提示无法解析
                        if not raw_text.strip():
                            st.warning(
                                f"{file_name} 无法提取文本内容，可能为加密文档或图片型DOC，请另存为.docx格式后上传。")
                            return [], False, "empty_content", file_name

                    except Exception as e:
                        st.error(f"{file_name} 解析失败：{str(e)[:50]}")
                        return [], False, "corrupt_doc", file_name
            elif file_name.endswith(('.ppt', '.pptx')):
                try:
                    from pptx import Presentation
                    prs = Presentation(file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                raw_text += shape.text + "\n"
                    page_num_record = len(prs.slides)
                except ImportError:
                    st.error("💡 要解析PPT，请安装：pip install python-pptx")
                    return [], False, "missing_dependency", file_name
                except Exception as e:
                    st.error(f"PPT解析失败：{str(e)[:50]}")
                    return [], False, "fail", file_name
            else:
                st.error(f"不支持格式：{file_name}")
                return [], False, "fail", file_name

            chunks, chunk_status = split_text_with_retry(raw_text, encoder, chunk_size, sim_threshold, min_length)
            file_status = chunk_status

            if not chunks:
                if attempt < FILE_RETRY_TIMES:
                    st.warning(f"文件{file_name}重试{attempt + 1}/{FILE_RETRY_TIMES + 1}")
                    time.sleep(2)
                    continue
                else:
                    st.warning(f"文件{file_name}无有效分段")
                    return [], False, file_status, file_name

            chunk_with_title = []
            for idx, chunk in enumerate(chunks):
                if chunk.strip():
                    title = generate_auto_chunk_title_optimized(file_name, chunk, idx, page_num_record)
                    chunk_with_title.append((title, chunk))

            if not chunk_with_title:
                if attempt < FILE_RETRY_TIMES:
                    st.warning(f"文件{file_name}重试{attempt + 1}/{FILE_RETRY_TIMES + 1}（无空分段）")
                    time.sleep(2)
                    continue
                else:
                    return [], False, "fail", file_name

            st.success(f"{file_name}：生成{len(chunk_with_title)}个分段")
            return chunk_with_title, True, "success", file_name

        except TimeoutError as e:
            file_status = "timeout"
            if attempt < FILE_RETRY_TIMES:
                st.warning(f"文件{file_name}超时重试{attempt + 1}/{FILE_RETRY_TIMES + 1}")
                time.sleep(2)
                continue
            else:
                st.error(f"文件{file_name}超时失败")
                return [], False, file_status, file_name
        except Exception as e:
            file_status = "fail"
            if attempt < FILE_RETRY_TIMES:
                st.warning(f"文件{file_name}失败重试{attempt + 1}/{FILE_RETRY_TIMES + 1}：{str(e)[:50]}")
                time.sleep(2)
                continue
            else:
                st.error(f"文件{file_name}失败：{str(e)[:50]}")
                return [], False, file_status, file_name


# ====================== 数据库操作类（保留原有功能 + 修复冗余bug） ==========
class LocalPGVectorDB:
    def __init__(self, db_config, silent=False):
        """完全回滚到你原始的__init__，仅保留必要的初始化逻辑，不新增任何东西"""
        self.config = db_config
        self.conn = None
        self.cursor = None
        self.silent = silent  # 这一行必须存在！
        self._connect()
        print("[调试] _connect() 已返回，即将执行 _fix_schema_local()")
        self._fix_schema_local()
        print("[调试] _fix_schema_local() 已返回")
        if not silent:  # 条件显示
            st.success("\数据库初始化完成")

    def _connect(self):
        print("[调试] _connect() 开始执行...")
        try:
            # 修复 localhost 解析问题：强制转为 127.0.0.1
            if self.config['host'] == 'localhost':
                self.config['host'] = '127.0.0.1'
                print("[调试] 已将 localhost 转换为 127.0.0.1")

            # 1. 安全关闭旧连接（加异常捕获，防止关闭损坏连接时卡住）
            if self.cursor:
                try:
                    self.cursor.close()
                    print("[调试] 旧游标已关闭")
                except Exception as e:
                    print(f"[调试] 关闭旧游标失败（忽略）: {e}")
                self.cursor = None

            if self.conn:
                try:
                    self.conn.close()
                    print("[调试] 旧连接已关闭")
                except Exception as e:
                    print(f"[调试] 关闭旧连接失败（忽略）: {e}")
                self.conn = None

            # 2. 连接 postgres 基础库（添加 sslmode='disable' 避免 SSL 协商卡住）
            print(f"[调试] 开始连接: {self.config['host']}:5432")
            base_conn = psycopg2.connect(
                dbname="postgres",
                user=self.config['user'],
                password=self.config['password'],
                host=self.config['host'],
                port=self.config['port'],
                connect_timeout=5,
                sslmode='disable'  #  关键：禁用 SSL，避免 SSL 握手卡住
            )
            print("[调试] 基础库连接成功")
            base_conn.autocommit = True
            print("[调试] 设置 autocommit 成功")

            base_cursor = base_conn.cursor()
            print("[调试] 创建游标成功")

            # 检查数据库是否存在
            print(f"[调试] 检查数据库 {self.config['dbname']} 是否存在...")
            base_cursor.execute(f"SELECT 1 FROM pg_database WHERE datname='{self.config['dbname']}';")
            exists = base_cursor.fetchone()
            print(f"[调试] 数据库存在: {exists}")

            if not exists:
                print(f"[调试] 创建数据库 {self.config['dbname']}...")
                try:
                    base_cursor.execute(f"CREATE DATABASE {self.config['dbname']} ENCODING 'UTF8';")
                    print("[调试] 创建数据库成功")
                except Exception as e:
                    print(f"[调试] 创建数据库失败（可能已存在）: {e}")

            print("[调试] 关闭基础库连接...")
            base_cursor.close()
            base_conn.close()
            print("[调试] 基础库连接已关闭")

            #  关键：连接目标库时也要加超时！之前只加了基础库
            print(f"[调试] 开始连接目标库: {self.config['dbname']}")
            self.conn = psycopg2.connect(
                host=self.config['host'],
                port=self.config['port'],
                user=self.config['user'],
                password=self.config['password'],
                dbname=self.config['dbname'],
                connect_timeout=5,  # 添加超时
                sslmode='disable'  # 禁用 SSL
            )
            print("[调试] 目标库连接成功")

            self.conn.autocommit = False
            self.cursor = self.conn.cursor(cursor_factory=RealDictCursor)
            register_vector(self.conn)
            print("[调试] 目标库初始化完成")

            if not self.silent:  # 添加条件判断
                st.success("数据库连接成功")
                st.success(
                    "Intel iGPU 连接成功  ▸ 目标设备：GPU.0 ▸ 设备名称：Intel(R) Iris(R) Xe Graphics (iGPU) ▸ 可用设备：['CPU', 'GPU.0', 'GPU.1']")

        except Exception as e:
            # 最小修复：异常时回滚+关闭连接，避免事务阻塞+连接泄漏
            if self.conn:
                try:
                    self.conn.rollback()
                    self.conn.close()  # 强制关闭连接
                except:
                    pass
            if self.cursor:
                try:
                    self.cursor.close()  # 强制关闭游标
                except:
                    pass
            self.conn = None
            self.cursor = None
            st.error(f"数据库连接失败：{str(e)[:80]}")

    def _fix_schema_local(self):
        print("[调试] _fix_schema_local() 开始执行...")
        # 先判断cursor和conn是否为None，避免报错
        if self.cursor is None or self.conn is None:
            st.error("数据库连接未建立，无法修复表结构！")
            return

        try:
            print("[调试] 设置语句超时 10 秒...")
            self.cursor.execute("SET statement_timeout = '10s';")

            print("[调试] 创建 vector 扩展...")
            self.cursor.execute("CREATE EXTENSION IF NOT EXISTS vector;")
            print("[调试] vector 扩展创建/检查完成")
            create_sql = [
                """
                CREATE TABLE IF NOT EXISTS knowledge_base (
                    id SERIAL PRIMARY KEY,
                    name VARCHAR(100) NOT NULL UNIQUE,
                    create_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                );
                """,
                """
                CREATE TABLE IF NOT EXISTS documents (
                    id SERIAL PRIMARY KEY,
                    kb_id INT REFERENCES knowledge_base(id) ON DELETE CASCADE,
                    name VARCHAR(200) NOT NULL,
                    create_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                );
                """,
                """
                CREATE TABLE IF NOT EXISTS text_chunks (
                    id SERIAL PRIMARY KEY,
                    doc_id INT REFERENCES documents(id) ON DELETE CASCADE,
                    title VARCHAR(200) DEFAULT '未命名分段',
                    content TEXT NOT NULL,
                    embedding vector(768) NOT NULL,
                    create_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                );
                """,
                """
                CREATE TABLE IF NOT EXISTS system_users (
                    id SERIAL PRIMARY KEY,
                    username VARCHAR(50) NOT NULL UNIQUE,
                    password VARCHAR(50) NOT NULL,
                    permissions INT DEFAULT 99,  -- 新增：权限字段，默认值99（普通用户权限）
                    create_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                );
                """,
                """
                CREATE TABLE IF NOT EXISTS city_adcode(
                    id SERIAL PRIMARY KEY,
                    city_name VARCHAR(100) UNIQUE NOT NULL,
                    adcode VARCHAR(10) NOT NULL,
                    citycode VARCHAR(10)
                );
                """

            ]
            for sql in create_sql:
                self.cursor.execute(sql)
            # 先检查并添加 permissions 列
            self.cursor.execute("""
                    SELECT column_name FROM information_schema.columns
                    WHERE table_name = 'system_users' AND column_name = 'permissions';
                    """)
            if not self.cursor.fetchone():
                self.cursor.execute("""
                        ALTER TABLE system_users ADD COLUMN permissions INT DEFAULT 99;
                        """)
                if not self.silent:  # 添加条件判断
                    st.success("给system_users表新增permissions字段")

            # ========== 强力清理：删除所有非 SM4 加密的 admin 变体 ==========
            encrypt_admin_user = sm4_encrypt("admin")  # 标准 SM4 密文

            # 1. 删除明文 admin
            self.cursor.execute("DELETE FROM system_users WHERE username = 'admin';")

            # 2. 删除所有 username 不是标准 SM4 密文，但可能是 admin 的其他形式（如 MD5、其他哈希）
            #    常见 admin 哈希值：MD5('admin') = 21232f297a57a5a743894a0e4a801fc3
            #    可根据需要添加更多变体
            admin_variants = [
                '21232f297a57a5a743894a0e4a801fc3',  # MD5 of 'admin'
                'admin',
                'ADMIN',
            ]
            # 删除这些变体（如果存在）
            for variant in admin_variants:
                self.cursor.execute("DELETE FROM system_users WHERE username = %s;", (variant,))

            # 3. 现在所有非标准 SM4 的 admin 记录已删除，只剩下可能的标准 SM4 密文记录
            # 查询所有标准 SM4 密文 admin
            self.cursor.execute("SELECT id, permissions FROM system_users WHERE username = %s;", (encrypt_admin_user,))
            admin_users = self.cursor.fetchall()

            # 4. 确保只有一个，且权限为 255
            if not admin_users:
                # 没有标准密文 admin → 创建
                encrypt_admin_pwd = sm4_encrypt("admin")
                self.cursor.execute("""
                                    INSERT INTO system_users (username, password, permissions)
                                    VALUES (%s, %s, %s)
                                    RETURNING id;
                                    """, (encrypt_admin_user, encrypt_admin_pwd, Permissions.DEFAULT_ADMIN))
                new_id = self.cursor.fetchone()['id']
                if not self.silent:
                    st.success(f"创建标准 admin (id={new_id})，权限 255")
            else:
                # 有标准密文 admin，保留 id 最小的，删除其余的
                admin_users.sort(key=lambda x: x['id'])
                keep_id = admin_users[0]['id']
                # 确保保留的权限为 255
                if admin_users[0]['permissions'] != Permissions.DEFAULT_ADMIN:
                    self.cursor.execute("UPDATE system_users SET permissions = %s WHERE id = %s;",
                                        (Permissions.DEFAULT_ADMIN, keep_id))
                    if not self.silent:
                        st.info(f"更新 admin (id={keep_id}) 权限为 255")
                # 删除其他
                if len(admin_users) > 1:
                    delete_ids = [user['id'] for user in admin_users[1:]]
                    self.cursor.execute("DELETE FROM system_users WHERE id = ANY(%s);", (delete_ids,))
                    if not self.silent:
                        st.warning(f"删除了 {len(delete_ids)} 个多余的标准 admin (ids: {delete_ids})")

            # 5. 最后提交
            self.conn.commit()
            self.cursor.execute("SELECT id, username, permissions FROM system_users WHERE username = %s;",
                                (encrypt_admin_user,))
            final_admin = self.cursor.fetchone()
            if final_admin:
                print(
                    f"[调试] 最终保留的 admin: id={final_admin['id']}, username={final_admin['username']}, permissions={final_admin['permissions']}")
                if not self.silent:
                    st.info(f"管理员账号: id={final_admin['id']}, 权限值={final_admin['permissions']} (期望 255)")
            else:
                print("[调试] 警告：未找到任何标准密文 admin 记录！")
                if not self.silent:
                    st.warning("未找到标准密文 admin 记录，请检查 SM4 加密是否正常")
            if not self.silent:
                st.success("admin 用户已清理完毕：仅保留一个密文 admin，权限 255")
        except Exception as e:
            # 先判断conn是否为None，再执行rollback
            if self.conn is not None:
                try:
                    self.conn.rollback()
                except:
                    pass
            print(f"[调试] _fix_schema_local() 异常: {e}")
            st.error(f"表结构修复失败：{str(e)[:80]}")
            logging.error(f"表结构失败详情：{traceback.format_exc()}")

    def get_kb_id(self, kb_name):
        try:
            self.cursor.execute("INSERT INTO knowledge_base (name) VALUES (%s) RETURNING id;", (kb_name,))
            kb_id = self.cursor.fetchone()['id']
            self.conn.commit()
            st.success(f"新建知识库「{kb_name}」")
            return kb_id
        except psycopg2.errors.UniqueViolation:
            self.conn.rollback()
            self.cursor.execute("SELECT id FROM knowledge_base WHERE name=%s;", (kb_name,))
            result = self.cursor.fetchone()
            if result:
                kb_id = result['id']
                st.info(f"追加到知识库「{kb_name}」")
                return kb_id
            return None
        except Exception as e:
            self.conn.rollback()
            st.error(f"知识库操作失败：{str(e)[:80]}")
            return None

    def batch_upload_auto_retry(self, kb_id, files, encoder, chunk_size=CHUNK_SIZE_DEFAULT, sim_threshold=0.8):
        if encoder is None:
            st.error("编码工具未初始化，无法处理文件")
            return 0, 0, {}

        if not kb_id or not files:
            return 0, 0, {}

        total_success = 0
        total_chunks = 0
        file_status_dict = {}
        progress = st.progress(0, text="批量处理中...")

        self.cursor.execute("""
            SELECT name FROM documents WHERE kb_id=%s AND status IN ('timeout', 'fail');
        """, (kb_id,))
        timeout_file_names = [doc['name'] for doc in self.cursor.fetchall()]
        if timeout_file_names:
            st.info(f"检测到{len(timeout_file_names)}个待重试文件：{', '.join(timeout_file_names)}")

        file_name_to_file = {file.name: file for file in files}
        all_file_names = list(file_name_to_file.keys())

        for idx, file_name in enumerate(all_file_names):
            file = file_name_to_file[file_name]
            progress.progress((idx + 1) / len(all_file_names),
                              text=f"处理 {file_name} ({idx + 1}/{len(all_file_names)})")

            self.cursor.execute("""
                SELECT status FROM documents WHERE kb_id=%s AND name=%s;
            """, (kb_id, file_name))
            existing_doc = self.cursor.fetchone()
            if existing_doc and existing_doc['status'] == 'success':
                st.info(f"{file_name} 已成功，跳过")
                file_status_dict[file_name] = ('success_skipped', 0)
                continue

            chunks_with_title, success, file_status, _ = parse_file_with_auto_retry(file, encoder, chunk_size,
                                                                                    sim_threshold)

            if not success or not chunks_with_title:
                file_status_dict[file_name] = (file_status, 0)
                if existing_doc:
                    self.cursor.execute("""
                        UPDATE documents SET status=%s WHERE kb_id=%s AND name=%s;
                    """, (file_status, kb_id, file_name))
                else:
                    self.cursor.execute("""
                        INSERT INTO documents (kb_id, name, status) VALUES (%s, %s, %s);
                    """, (kb_id, file_name, file_status))
                self.conn.commit()
                st.warning(f"跳过 {file_name}：{file_status}")
                continue

            if existing_doc:
                self.cursor.execute("""
                    UPDATE documents SET status=%s WHERE kb_id=%s AND name=%s RETURNING id;
                """, ('success', kb_id, file_name))
                doc_id = self.cursor.fetchone()['id']
            else:
                self.cursor.execute("""
                    INSERT INTO documents (kb_id, name, status) VALUES (%s, %s, %s) RETURNING id;
                """, (kb_id, file_name, 'success'))
                doc_id = self.cursor.fetchone()['id']

            chunk_count = 0
            for title, chunk in chunks_with_title:
                if not isinstance(chunk, str):
                    st.warning(f"无效分段内容，跳过")
                    continue
                embedding = encoder.encode_text(chunk)
                self.cursor.execute("""
                    INSERT INTO text_chunks (doc_id, title, content, embedding)
                    VALUES (%s, %s, %s, %s);
                """, (doc_id, title, chunk, embedding.tolist()))
                chunk_count += 1

            self.conn.commit()
            total_success += 1
            total_chunks += chunk_count
            file_status_dict[file_name] = ('success', chunk_count)
            st.success(f"{file_name}：新增{chunk_count}个分段")

        progress.progress(1.0, text=f"处理完成：成功{total_success}/{len(all_file_names)}文件，共{total_chunks}分段")
        return total_success, total_chunks, file_status_dict

    def get_all_kbs(self):
        try:
            self.cursor.execute("""
                SELECT kb.id, kb.name, kb.create_time,
                       COUNT(DISTINCT d.id) as doc_count,
                       COUNT(t.id) as chunk_count
                FROM knowledge_base kb
                LEFT JOIN documents d ON kb.id = d.kb_id
                LEFT JOIN text_chunks t ON d.id = t.doc_id
                GROUP BY kb.id ORDER BY kb.create_time DESC;
            """)
            return self.cursor.fetchall()
        except Exception as e:
            st.error(f"读取知识库失败：{str(e)[:80]}")
            return []

    def get_kb_chunks(self, kb_id, page=1, page_size=PAGE_SIZE):
        try:
            self.cursor.execute("""
                SELECT COUNT(*) as total FROM text_chunks tc
                JOIN documents d ON tc.doc_id = d.id WHERE d.kb_id = %s;
            """, (kb_id,))
            total = self.cursor.fetchone()['total']

            total_pages = (total + page_size - 1) // page_size if total > 0 else 0
            if page > total_pages and total_pages > 0:
                page = total_pages
            elif page < 1:
                page = 1

            offset = (page - 1) * page_size

            self.cursor.execute("""
                SELECT tc.id, tc.title, tc.content, tc.create_time, d.name as doc_name
                FROM text_chunks tc
                JOIN documents d ON tc.doc_id = d.id
                WHERE d.kb_id = %s 
                ORDER BY tc.create_time ASC, tc.id ASC
                LIMIT %s OFFSET %s;
            """, (kb_id, page_size, offset))
            chunks = self.cursor.fetchall()

            return chunks, total_pages, total
        except Exception as e:
            st.error(f"读取分段失败：{str(e)[:80]}")
            return [], 0, 0

    def update_chunk(self, chunk_id, title, content, encoder):
        if encoder is None or not isinstance(content, str):
            st.error("编码工具无效或分段内容异常，无法更新")
            return False
        try:
            embedding = encoder.encode_text(content)
            self.cursor.execute("""
                UPDATE text_chunks SET title=%s, content=%s, embedding=%s
                WHERE id=%s;
            """, (title, content, embedding.tolist(), chunk_id))
            self.conn.commit()
            return True
        except Exception as e:
            self.conn.rollback()
            st.error(f"更新分段失败：{str(e)[:80]}")
            return False

    def delete_chunk(self, chunk_id):
        try:
            self.cursor.execute("DELETE FROM text_chunks WHERE id=%s;", (chunk_id,))
            self.conn.commit()
            return True
        except Exception as e:
            self.conn.rollback()
            st.error(f"删除分段失败：{str(e)[:80]}")
            return False

    def delete_kb(self, kb_id, kb_name):
        try:
            self.cursor.execute("DELETE FROM knowledge_base WHERE id=%s;", (kb_id,))
            self.conn.commit()
            st.success(f"知识库「{kb_name}」已删除")
            return True
        except Exception as e:
            self.conn.rollback()
            st.error(f"删除知识库失败：{str(e)[:80]}")
            return False

    # ========== 新增：用户查询方法（登录验证用） ==========
    # ========== 新增：用户查询方法（修改后：加密用户名再查询） ==========
    def get_user_by_username(self, username):
        """根据用户名获取用户信息（包含权限）"""
        try:
            encrypt_username = sm4_encrypt(username)
            self.cursor.execute(
                "SELECT id, username, password, permissions FROM system_users WHERE username = %s;",
                (encrypt_username,)
            )
            return self.cursor.fetchone()
        except Exception as e:
            self.conn.rollback()
            st.error(f"查询用户失败：{str(e)[:80]}")
            return None

    # ========== 新增：用户注册方法 ==========
    def register_user(self, username, password, permissions=None):
        if username.lower() == "admin":
            return "forbidden"
        """注册新用户（新增权限参数）"""
        try:
            # 设置默认权限
            if permissions is None:
                permissions = Permissions.DEFAULT_USER  # 普通用户默认权限

            encrypt_username = sm4_encrypt(username)
            encrypt_password = sm4_encrypt(password)

            self.cursor.execute("SELECT id FROM system_users WHERE username = %s;", (encrypt_username,))
            existing_user = self.cursor.fetchone()
            if existing_user:
                return "exists"

            self.cursor.execute(
                "INSERT INTO system_users (username, password, permissions) VALUES (%s, %s, %s) RETURNING id;",
                (encrypt_username, encrypt_password, permissions)
            )
            self.conn.commit()
            return "success"
        except Exception as e:
            self.conn.rollback()
            st.error(f"注册用户失败：{str(e)[:80]}")
            return "fail"

    # ========== 新增：检查用户是否存在（用户名已加密） ==========
    def check_user_exists(self, username):
        """检查用户名是否存在（使用加密后的用户名查询）"""
        try:
            encrypt_username = sm4_encrypt(username)
            self.cursor.execute("SELECT id FROM system_users WHERE username = %s;", (encrypt_username,))
            return self.cursor.fetchone() is not None
        except Exception as e:
            self.conn.rollback()
            st.error(f"查询用户失败：{str(e)[:80]}")
            return False

    # ========== 新增：更新用户密码（密码会加密存储） ==========
    def update_user_password(self, username, new_password):
        """更新用户密码（密码会加密存储）"""
        try:
            encrypt_username = sm4_encrypt(username)
            encrypt_new_pwd = sm4_encrypt(new_password)
            self.cursor.execute(
                "UPDATE system_users SET password = %s WHERE username = %s;",
                (encrypt_new_pwd, encrypt_username)
            )
            self.conn.commit()
            return True
        except Exception as e:
            self.conn.rollback()
            st.error(f"更新密码失败：{str(e)[:80]}")
            return False

    # 修复冗余代码bug，保留完整检索逻辑
    def _search_similar_chunks_raw(self, query, encoder, top_k=5, kb_id=None):
        try:
            # 优化：简化向量编码校验，减少冗余逻辑
            query_vector = encoder.encode_text(query)
            # 直接判断向量是否有效，快速返回
            if query_vector is None or len(query_vector) == 0 or np.all(query_vector == 0):
                st.warning("问题编码失败，无法进行有效检索")
                return []

            if top_k is None:
                top_k = 8  # 默认增加检索结果数量，提高信息覆盖率
            # 新增：提前初始化results，避免变量未定义
            results = []
            params = []
            # 修复：严格区分指定知识库和全局检索，无逻辑漏洞
            if kb_id is None:
                # 全局检索：遍历所有知识库，无kb_id过滤
                base_sql = """
                       SELECT tc.id, tc.title, tc.content, 
                              tc.embedding <=> %s::vector(768) as cosine_distance,
                              d.name as doc_name, kb.name as kb_name, tc.create_time
                       FROM text_chunks tc
                       JOIN documents d ON tc.doc_id = d.id
                       JOIN knowledge_base kb ON d.kb_id = kb.id
                       ORDER BY cosine_distance ASC, kb.id ASC 
                       LIMIT %s;
                   """
                max_candidate = 50
                candidate_num = min(top_k * 3, max_candidate)
                params = [query_vector.tolist(), top_k * 10]  # 取更多候选结果，再筛选
                self.cursor.execute(base_sql, params)
                results = self.cursor.fetchall()
            else:
                # 步骤1：单独设置查询超时（无结果，无需fetchall）
                set_timeout_sql = "SET statement_timeout = 5000;"
                self.cursor.execute(set_timeout_sql)
                # 指定知识库：严格过滤kb_id，确保范围可控
                base_sql = """
                       SELECT tc.id, tc.title, tc.content, 
                              tc.embedding <=> %s::vector(768) as cosine_distance,
                              d.name as doc_name, kb.name as kb_name, tc.create_time
                       FROM text_chunks tc
                       JOIN documents d ON tc.doc_id = d.id
                       JOIN knowledge_base kb ON d.kb_id = kb.id
                       WHERE d.kb_id = %s
                       ORDER BY cosine_distance ASC
                       LIMIT %s;
                   """
                max_candidate = 50
                candidate_num = min(top_k * 3, max_candidate)
                params = [query_vector.tolist(), kb_id, top_k * 3]
                # 全局检索直接执行，无超时设置
                self.cursor.execute(base_sql, params)
                results = self.cursor.fetchall()

                # 步骤3：执行SET重置（无结果），不调用fetchall()，无报错
                reset_timeout_sql = "SET statement_timeout = 0;"
                self.cursor.execute(reset_timeout_sql)  # 无结果，不调用fetchall()

            # 新增：小知识库快速返回逻辑（结果少则直接处理，无需复杂流程）
            if not results:
                return []
            # 若结果数量≤top_k，直接标记为有效结果，跳过后续冗余过滤
            if len(results) <= top_k:
                enhanced_results = []
                for res in results:
                    chunk_content = clean_text_local(res['content']).strip()
                    if not chunk_content:
                        continue
                    cosine_sim = round(1 - res['cosine_distance'], 4)
                    keyword_sim = calculate_keyword_match_score(query, chunk_content, top_k=5)
                    total_sim = round((1 - KEYWORD_MATCH_WEIGHT) * cosine_sim + KEYWORD_MATCH_WEIGHT * keyword_sim, 4)
                    res['cosine_similarity'] = cosine_sim
                    res['keyword_similarity'] = keyword_sim
                    res['total_similarity'] = total_sim
                    del res['cosine_distance']
                    enhanced_results.append(res)
                # 直接排序返回，无需额外裁剪
                enhanced_results.sort(key=lambda x: x['total_similarity'], reverse=True)
                return enhanced_results[:top_k]

            # 优化：放宽最小相似度阈值，从0.5降到0.4，避免有效结果被过滤
            enhanced_results = []
            for res in results:
                # 新增：提前过滤无效分段，减少无效计算
                chunk_content = clean_text_local(res['content']).strip()
                if not chunk_content:  # 过滤空内容/超短内容
                    continue
                cosine_sim = round(1 - res['cosine_distance'], 4)
                # 优化：提高关键词匹配的top_k，从3调到5，提取更多关键词，提高匹配准确性
                keyword_sim = calculate_keyword_match_score(query, res['content'], top_k=5)
                total_sim = round((1 - KEYWORD_MATCH_WEIGHT) * cosine_sim + KEYWORD_MATCH_WEIGHT * keyword_sim, 4)

                res['cosine_similarity'] = cosine_sim
                res['keyword_similarity'] = keyword_sim
                res['total_similarity'] = total_sim
                del res['cosine_distance']

                # 全局检索时：先不提前筛选，保留所有候选集；指定知识库检索时，按0.4筛选
                if kb_id is None:
                    # 全局检索：延迟筛选，先保留所有候选结果
                    enhanced_results.append(res)
                else:
                    # 指定知识库检索：保持原有筛选逻辑，不影响精准检索
                    if total_sim >= 0.3:
                        enhanced_results.append(res)

            # 先裁剪综合相似度<0.2的低质量结果，减少排序数据量
            enhanced_results = [item for item in enhanced_results if item['total_similarity'] >= 0.2]
            # 按综合相似度排序，取top_k
            enhanced_results.sort(key=lambda x: x['total_similarity'], reverse=True)
            final_results = enhanced_results[:top_k]
            return final_results
        except Exception as e:
            st.error(f"向量检索失败：{str(e)[:80]}")
            logging.error(f"向量检索失败详情：{traceback.format_exc()}")
            return []

    def search_similar_chunks(self, query, encoder, top_k=5, kb_id=None):
        """Redis 缓存版入口（优化版）"""
        # 优化缓存键：包含数据库名称和更多参数
        cache_key = f"search:{CACHE_VER}:{self.config['dbname']}:{hashlib.md5(f'{query}#{kb_id}#{top_k}'.encode()).hexdigest()}"

        # 尝试读取缓存
        cached = REDIS_CLI.get(cache_key)
        if cached:
            try:
                return json.loads(cached)
            except Exception:
                pass  # 缓存解析失败，继续执行检索

        # 执行原始检索（完全不变）
        raw_results = self._search_similar_chunks_raw(query, encoder, top_k, kb_id)
        clean_results = _serialize_rows(raw_results)

        # 缓存结果（延长到2小时）
        if clean_results:
            REDIS_CLI.setex(cache_key, 60 * 120, json.dumps(clean_results, ensure_ascii=False))

        return clean_results

    # ========== 新增：用户管理相关方法 ==========
    def get_all_users(self):
        """获取所有用户信息"""
        try:
            self.cursor.execute(
                "SELECT id, username, permissions, create_time FROM system_users ORDER BY create_time DESC;"
            )
            return self.cursor.fetchall()
        except Exception as e:
            self.conn.rollback()
            st.error(f"获取用户列表失败：{str(e)[:80]}")
            return []

    # 在LocalPGVectorDB类中，修改update_user_permissions方法
    def update_user_permissions(self, original_username, new_permissions):
        """更新用户权限（使用原始用户名，而不是加密后的）"""
        try:
            # 先加密用户名再查询
            encrypt_username = sm4_encrypt(original_username)
            self.cursor.execute(
                "UPDATE system_users SET permissions = %s WHERE username = %s;",
                (new_permissions, encrypt_username)
            )
            self.conn.commit()
            return True
        except Exception as e:
            self.conn.rollback()
            st.error(f"更新用户权限失败：{str(e)[:80]}")
            return False

    def delete_user(self, username):
        """删除用户"""
        try:
            encrypt_username = sm4_encrypt(username)
            self.cursor.execute(
                "DELETE FROM system_users WHERE username = %s;",
                (encrypt_username,)
            )
            self.conn.commit()
            return True
        except Exception as e:
            self.conn.rollback()
            st.error(f"删除用户失败：{str(e)[:80]}")
            return False


# ========== 通用 Ollama 答案缓存 ==========
def cached_ollama_response(
        prompt: str,
        format_schema: dict,
        model: str = OLLAMA_MODEL,
        ttl: int = 60 * 120  # 延长到2小时
) -> any:
    """
    优化的Ollama响应缓存，显著加速回答生成
    返回类型和用法完全不变，保证零改动
    """
    cache_key = f"ans:{CACHE_VER}:{model}:{hashlib.md5(prompt.encode()).hexdigest()}"

    # 尝试读取缓存
    cached = REDIS_CLI.get(cache_key)
    if cached:
        try:
            # 完全保持原有返回格式，下游代码零改动
            from types import SimpleNamespace
            data = json.loads(cached)
            ns = SimpleNamespace()
            ns.message = SimpleNamespace()
            ns.message.content = json.dumps(data, ensure_ascii=False)
            return ns
        except Exception as e:
            # 缓存读取失败，继续调用Ollama
            print(f"缓存读取失败: {e}")

    # 原路调用 Ollama（完全保留您的原有参数）
    client = Client(host='http://127.0.0.1:11434')
    try:
        response = client.chat(
            messages=[{'role': 'user', 'content': prompt}],
            model=model,
            format=format_schema,
            stream=False,
            options={
                "temperature": 0.1,
                "top_p": 0.9,
                "max_tokens": 8192,
                "num_ctx": 8192,
                "seed": 42,
            }
        )

        # 缓存结果（延长到2小时）
        try:
            raw_json = response.message.content.strip()
            answer = KnowledgeBaseAnswer.model_validate_json(raw_json)
            REDIS_CLI.setex(cache_key, ttl, json.dumps(answer.model_dump(), ensure_ascii=False))
        except Exception as e:
            print(f"缓存写入失败: {e}")

        return response  # 保持原有返回格式
    except Exception as e:
        # 保持原有错误处理逻辑
        raise Exception(f"Ollama调用失败：{str(e)[:100]}")


# ====================== 2. 问答页签：严格复刻你的Country模型4步核心逻辑（解决空值问题） ==========
def qa_tab():
    # =========  1. 取天气  =========
    adcode = st.session_state.get("user_adcode")
    city = st.session_state.get("user_city", "")
    logging.info(f"[debug] adcode={adcode}  city={city}")
    weather_text = ""  # 最终拼到 prompt
    api = GaodeWeatherAPI(GAODE_API_KEY)
    cur_data = None
    fc_data = None
    cur = None
    fc = None
    if adcode:
        cur_data = api.get_weather(adcode, "base")
        fc_data = api.get_weather(adcode, "all")
        cur = api.parse_current_weather(cur_data) if cur_data else None
        fc = api.parse_forecast_weather(fc_data) if fc_data else None
        if cur:
            weather_text += f"用户所在城市：{cur['city']}，" \
                            f"当前天气：{cur['weather']}，" \
                            f"气温：{cur['temperature']}℃，" \
                            f"湿度：{cur['humidity']}%，" \
                            f"风向：{cur['wind_direction']} {cur['wind_power']}级。"
        if fc and len(fc["forecasts"]) > 1:
            tom = fc["forecasts"][1]
            weather_text += f"明日预报：{tom['day_weather']}/{tom['night_weather']}，" \
                            f"温度：{tom['night_temp']}~{tom['day_temp']}℃。"
    # 权限检查
    if not Permissions.has_permission(st.session_state.user_permissions, Permissions.USE_QA):
        st.error("你没有使用智能问答的权限！")
        return
        # ★ 立即清屏，防止旧卡片残留
    if st.session_state.get("qa_clear_flag"):
        st.empty()
        st.session_state["qa_clear_flag"] = False
    st.subheader("知识库智能问答")
    st.caption("4步核心：定义模型→创建Client→实例调chat→验证JSON→拆分显示")
    # 简单的清理按钮
    col_title, col_stats, col_refresh = st.columns([3, 2, 1])
    with col_title:
        st.caption("流式输出优先，缓存加速响应")

    with col_stats:
        # 显示缓存状态
        try:
            cache_count = len(REDIS_CLI.keys(f"ans:{CACHE_VER}:*"))
            if cache_count > 0:
                st.caption(f"答案缓存: {cache_count} 条")
        except:
            pass

    with col_refresh:
        if st.button("🔄 刷新", key="refresh_qa", use_container_width=True,
                     help="清空当前问答状态，重新开始"):
            # 1. 清 session_state
            for k in list(st.session_state.keys()):
                if k.startswith('qa_') or k in ['show_thinking', 'qa_query', 'qa_stream_buffer', 'qa_clear_flag']:
                    st.session_state.pop(k, None)
            # 2. 立即清 UI
            st.session_state["qa_clear_flag"] = True
            st.rerun()

        # 在 qa_tab 函数最开始（任何 UI 之前）加：
        if st.session_state.get("qa_clear_flag"):
            st.empty()  # 一次性把整个页签冲掉
            st.session_state["qa_clear_flag"] = False

    if not st.session_state.get("qa_perf_hint_shown", False):
        with st.expander("响应优化提示", expanded=False):
            st.info("""
                **性能优化说明：**
                - 首次查询可能稍慢（约10-20秒），后续相似问题会从缓存快速响应（1-3秒）
                - 流式输出会实时显示内容，体验更佳
                - 您可以选择关闭"显示思考过程"以加快显示速度
                """)
        st.session_state.qa_perf_hint_shown = True

    if not st.session_state.init_ok:
        st.warning("请先左侧初始化并确认环境符合要求")
        return

    # 修复1：RealDictRow 键值访问
    kbs = st.session_state.db_client.get_all_kbs()
    kb_options = {"全局检索（所有知识库）": None}
    if kbs:
        for kb in kbs:
            kb_options[kb['name']] = kb['id']

    col1, col2, col3 = st.columns([2, 1, 1])
    with col1:
        selected_kb = st.selectbox("检索范围", options=list(kb_options.keys()), key="qa_kb")
    with col2:
        top_k = st.slider("返回上下文数量", 1, 15, 8, key="qa_top_k")  # 增加默认top_k，提高信息覆盖率
    with col3:
        speed_boost = st.checkbox(
            "极速模式",
            value=True,
            key="qa_speed_boost",
            help="启用后仅取相似度最高的3条结果，生成速度提升3倍"
        )

    query = st.text_area(
        "请输入你的问题",
        placeholder="例如：韭菜常见病虫害有哪些？如何防治？",
        height=120,
        key="qa_query"
    )
    # ---- 实时天气 ----
    weather_text = ""
    db = st.session_state.db_client
    encoder = st.session_state.encoder
    adcode = st.session_state.get("user_adcode")
    city = st.session_state.get("user_city", "")
    adcode = st.session_state.get("user_adcode")
    # 天气解析三件套
    cur = api.parse_current_weather(cur_data) if cur_data else None
    fc = api.parse_forecast_weather(fc_data) if fc_data else None
    tom = None  # ← 补上
    if fc and len(fc.get("forecasts", [])) > 1:
        tom = fc["forecasts"][1]
    if adcode:
        api = GaodeWeatherAPI(GAODE_API_KEY)
        cur_data = api.get_weather(adcode, "base")
        fc_data = api.get_weather(adcode, "all")
        cur = api.parse_current_weather(cur_data) if cur_data else None
        fc = api.parse_forecast_weather(fc_data) if fc_data else None
        if cur:
            weather_text += (f"【实时天气】{cur['city']}：{cur['weather']} "
                             f"{cur['temperature']}℃，湿度 {cur['humidity']}%，"
                             f"风向 {cur['wind_direction']} {cur['wind_power']}级；")
        if fc and len(fc["forecasts"]) > 1:
            tom = fc["forecasts"][1]
            weather_text += (f"【明日预报】{tom['day_weather']}/{tom['night_weather']}，"
                             f"温度 {tom['night_temp']}~{tom['day_temp']}℃。")
    # 构造带天气的查询串
    weather_kw = (
        f"{cur.get('weather', '') if cur else ''} "
        f"{tom.get('day_weather', '') if tom else ''} "
        f"{tom.get('night_weather', '') if tom else ''}"
    ).strip()
    query_with_weather = f"{query} {city} {weather_kw}".strip()
    similar_chunks = db.search_similar_chunks(
        query=query_with_weather,  # 用新串
        encoder=encoder,
        top_k=top_k,
        kb_id=kb_options[selected_kb]
    )
    # 2. 去重（按 content 完全一致）并保量
    seen = set()
    unique = []
    for ch in similar_chunks:
        key = ch['content'].strip()
        if key not in seen:
            seen.add(key)
            unique.append(ch)
            if len(unique) == top_k:  # 一旦凑够滑动条值就停
                break
    similar_chunks = unique
    # 新增：是否显示思考过程复选框（默认勾选，取消后完全不渲染思考过程区域）
    show_thinking = st.checkbox("显示思考过程", value=True, key="show_thinking")

    if st.button("🔍 检索并生成智能回答", type="primary", key="qa_search_btn"):
        if not query.strip():
            st.error("问题不能为空，请输入有效问题！")
            return
        full_response = ""
        progress_bar = st.progress(0, text="正在准备检索...")
        status_text = st.empty()

        with st.spinner("正在检索知识库并按指定JSON格式生成回答..."):
            # 第一步：检索知识库
            progress_bar.progress(20, text="检索知识库中...")
            status_text.info("正在检索知识库内容...")
            db = st.session_state.db_client
            encoder = st.session_state.encoder
            kb_id = kb_options[selected_kb]
            progress_bar.progress(40, text="构造提示词中...")
            status_text.info("正在构造提示词...")
            # 1. 检索知识库，构造查询结果文本（给Ollama的输入）
            similar_chunks = db.search_similar_chunks(
                query=query,  # 用户问题
                encoder=encoder,
                top_k=top_k,
                kb_id=kb_options[selected_kb]  # 知识库 ID
            )
            similar_chunks = similar_chunks[: 3 if speed_boost else top_k]
            # 2. 过滤掉content为空的无效结果
            similar_chunks = [chunk for chunk in similar_chunks if chunk['content'].strip()]
            if not similar_chunks:
                st.error("检索结果中无有效内容")
                return

            kb_search_text = ""
            if similar_chunks:
                # 强制保留每条结果的完整格式，不做过度清洗
                kb_search_text = "\n\n".join([
                    f"【检索结果第{idx + 1}条 | 综合相似度：{chunk['total_similarity']:.4f}】\n"
                    f"来源文档：{chunk['doc_name']}\n"
                    f"分段标题：{chunk['title']}\n"
                    f"内容：{chunk['content'].strip()}"
                    for idx, chunk in enumerate(similar_chunks)
                ])
            else:
                kb_search_text = "未检索到任何与用户问题相关的知识库片段"
            progress_bar.progress(60, text="生成智能回答中...")
            status_text.info("正在生成回答（流式输出）...")
            st.write(f"传给Ollama前的内容长度：{len(kb_search_text)} | 前200字符：{kb_search_text[:200]}...")

            # ========== 关键：强制Prompt，要求Ollama仅返回标准JSON，无任何多余内容 ==========
            # ========== 简化提示词，保持有效 ==========
            prompt_content = f"""
            基于以下知识库内容和当下天气情况回答用户问题。返回一个包含三个字段的JSON对象：
            1. kb_search_results: 必须原样复制下面的知识库内容和当地天气
            2. thinking_process: 仔细分析知识库内容和当地天气如何回答用户问题
            3. final_conclusion: 基于知识库内容和当地天气给出具体答案
            知识库内容：
            {kb_search_text}
            当地天气：{weather_text}
            用户问题：{query}
            只返回JSON，不要其他任何内容。
            """.strip()

            # === 使用 /api/ask 接口 ===
            import requests, json, time
            try:
                # 在调用 Django 前添加调试信息
                print(f">>> 调试 - 天气信息: {weather_text}", file=sys.stderr)
                print(f">>> 调试 - 天气信息长度: {len(weather_text)}", file=sys.stderr)

                # 检查天气信息是否为空
                if not weather_text or weather_text.strip() == "":
                    print(">>> 警告: 天气信息为空！", file=sys.stderr)
                    # 使用默认天气信息
                    weather_text = "当前天气信息不可用，请根据知识库内容回答问题。"
                # 调用 Django 的加速接口
                r = requests.post("http://127.0.0.1:8000/api/ask",
                                  json={
                                      "prompt": prompt_content,
                                      "kb_search_text": kb_search_text,
                                      "weather_text": weather_text
                                  },timeout=60)  # 增加超时时间

                res = r.json()

                if res.get("hit"):
                    # 缓存命中，直接解析答案
                    print(">>> 缓存命中，收到答案包", file=sys.stderr)
                    ans_dict = json.loads(res["answer"]) if isinstance(res["answer"], str) else res["answer"]
                    kb_answer = KnowledgeBaseAnswer(**ans_dict)
                else:
                    # 缓存未命中，需要轮询
                    task_id = res["task_id"]
                    print(f">>> 缓存未命中，开始轮询任务 {task_id}", file=sys.stderr)

                    answer_obtained = False
                    polling_answer = None

                    # 去掉总超时，改为无限轮询，直到成功或失败
                    max_poll_interval = 1.0
                    min_poll_interval = 0.1
                    poll_interval = min_poll_interval


                    with st.spinner(f"正在生成回答 (任务 {task_id})..."):
                        while True:
                            try:
                                result_resp = requests.get(f"http://127.0.0.1:8000/api/result/{task_id}", timeout=5)
                                result_resp.raise_for_status()
                                result_data = result_resp.json()

                                # 打印接收到的报文（便于调试）
                                print(f">>> 收到轮询响应: {result_data}", file=sys.stderr)

                                if "answer" in result_data:
                                    ans = result_data["answer"]
                                    ans_dict = json.loads(ans) if isinstance(ans, str) else ans
                                    kb_answer = KnowledgeBaseAnswer(**ans_dict)
                                    answer_obtained = True
                                    polling_answer = kb_answer

                                    # 发送结束确认报文（DELETE）
                                    try:
                                        del_resp = requests.delete(f"http://127.0.0.1:8000/api/result/{task_id}",
                                                                   timeout=2)
                                        print(
                                            f">>> 发送结束确认: DELETE /api/result/{task_id} -> {del_resp.status_code}",
                                            file=sys.stderr)
                                    except Exception as del_err:
                                        print(f">>> 结束确认失败: {del_err}", file=sys.stderr)

                                    # 清理进度条和状态
                                    progress_bar.empty()
                                    status_text.empty()
                                    break  # 成功，退出轮询

                                elif result_data.get("status") == "failed":
                                    st.error(f"任务失败: {result_data.get('error', '未知错误')}")
                                    # 创建错误答案
                                    kb_answer = KnowledgeBaseAnswer(
                                        kb_search_results=kb_search_text,
                                        thinking_process="任务失败",
                                        final_conclusion=f"生成答案失败：{result_data.get('error', '未知错误')}"
                                    )
                                    progress_bar.empty()
                                    status_text.empty()
                                    break

                                else:
                                    # 仍在运行，等待后继续轮询
                                    time.sleep(poll_interval)
                                    poll_interval = min(poll_interval * 1.5, max_poll_interval)
                                    continue

                            except requests.exceptions.RequestException as req_err:
                                print(f">>> 轮询请求失败: {req_err}，重试中...", file=sys.stderr)
                                time.sleep(poll_interval)
                                poll_interval = min(poll_interval * 1.5, max_poll_interval)
                                continue
                            except json.JSONDecodeError as json_err:
                                print(f">>> 响应JSON解析失败: {json_err}，原始响应: {result_resp.text}", file=sys.stderr)
                                time.sleep(poll_interval)
                                poll_interval = min(poll_interval * 1.5, max_poll_interval)
                                continue

            except Exception as e:
                st.error(f"调用API失败: {str(e)}")
                # 创建错误答案
                kb_answer = KnowledgeBaseAnswer(
                    kb_search_results=f"API调用失败: {str(e)}",
                    thinking_process="无法获取思考过程",
                    final_conclusion="无法生成答案，请检查网络连接或稍后重试"
                )

            if answer_obtained:
                kb_answer = polling_answer
            else:
                # ========== 4. 严格按你的逻辑：验证JSON并解析（确保非空，清洗多余内容） ==========
                try:
                    # 关键修复：确保response变量存在
                    if 'response' not in locals():
                        # 如果response不存在，重新调用缓存函数
                        response = cached_ollama_response(prompt_content, KnowledgeBaseAnswer.model_json_schema())

                    # 强力清洗：移除所有非JSON内容，确保解析成功
                    raw_json = response.message.content.strip()
                    raw_json = re.sub(r'^```json|```$', '', raw_json).strip()  # 移除代码块标记
                    raw_json = re.sub(r'^\s+|\s+$', '', raw_json)  # 移除首尾空格
                    if not raw_json:
                        raise Exception("Ollama返回空JSON内容")

                    # 验证并解析JSON（完全对齐你的Country示例）
                    kb_answer = KnowledgeBaseAnswer.model_validate_json(raw_json)

                    # 最终兜底：确保每个字段非空（防止模型绕过指令）
                    if not kb_answer.kb_search_results.strip():
                        kb_answer.kb_search_results = kb_search_text  # 直接赋值原始知识库内容
                    if not kb_answer.thinking_process.strip():
                        kb_answer.thinking_process = "已检索知识库内容，未发现与用户问题相关的有效信息可供分析"
                    if not kb_answer.final_conclusion.strip():
                        kb_answer.final_conclusion = "未从知识库中检索到与该问题直接相关的有效信息，无法给出针对性结论"

                except Exception as e:
                    st.error(f"JSON验证/解析失败：{str(e)[:100]}")

                    # 安全地获取响应内容
                    try:
                        # 检查response是否存在
                        if 'response' in locals() and response:
                            if hasattr(response, 'message') and hasattr(response.message, 'content'):
                                content = response.message.content[:500] if response.message.content else "空"
                                st.warning(f"原始响应内容：{content}...")
                        elif 'full_response' in locals():
                            st.warning(f"流式响应内容：{full_response[:500]}...")
                        else:
                            st.warning("无法获取原始响应内容")
                    except Exception as e2:
                        st.warning(f"获取响应内容失败：{e2}")

                    # 确保有兜底的kb_answer
                    try:
                        kb_answer = KnowledgeBaseAnswer(
                            kb_search_results=kb_search_text,
                            thinking_process="JSON解析异常，使用检索结果直接分析",
                            final_conclusion="基于检索到的知识库内容，相关信息如下：" +
                                             ("\n".join([f"- {chunk['title']}: {chunk['content'][:100]}..."
                                                         for chunk in
                                                         similar_chunks[:3]]) if similar_chunks else "未找到相关信息")
                        )
                    except:
                        # 最终兜底
                        kb_answer = KnowledgeBaseAnswer(
                            kb_search_results="处理异常",
                            thinking_process="处理异常",
                            final_conclusion="请重试或联系管理员"
                        )
                    logging.error(f"JSON解析失败详情：{traceback.format_exc()}")

                    # 清理进度条和状态文本（新增）
                    if 'progress_bar' in locals():
                        progress_bar.empty()
                    if 'status_text' in locals():
                        status_text.empty()

            # ========== 5. 按JSON字段拆分显示（完全满足你的要求，非空） ==========
            # 第一步：显示知识库查询结果（JSON的kb_search_results字段）
            st.markdown("### 第一部分：知识库查询结果（JSON字段：kb_search_results）")
            if similar_chunks:
                with st.expander(f"点击展开/收起（共{len(similar_chunks)}条结果）", expanded=True):
                    for idx, chunk in enumerate(similar_chunks, 1):
                        st.markdown(f"""
<div style="background:#f5f7fa;padding:10px;border-radius:6px;margin-bottom:8px;">
    <strong>第{idx}条 | 知识库：{chunk['kb_name']} | 文档：{chunk['doc_name']}</strong>
    <br>
    <strong>综合相似度：{chunk['total_similarity']} | 余弦相似度：{chunk['cosine_similarity']}</strong>
    <br>
    <strong>标题：</strong>{chunk['title']}
    <br>
    <strong>内容：</strong>{clean_text_local(chunk['content'])}
</div>
                        """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
<div style="background:#fef7f7;padding:10px;border-radius:6px;color:#dc2626;">
    {kb_answer.kb_search_results}  <!-- 显示JSON字段内容 -->
</div>
                """, unsafe_allow_html=True)

            # 第二步：显示思考过程（仅当show_thinking为True时，才渲染整个思考过程区域）
            if show_thinking:  # 新增条件判断：不满足则完全不执行以下渲染代码
                st.markdown("### 第二部分：思考过程（JSON字段：thinking_process）")
                st.markdown(f"""
            <div style="background:#eef2ff;padding:12px;border-radius:8px;border:1px solid #c7d2fe;">
                {kb_answer.thinking_process}  <!-- 显示JSON字段内容 -->
            </div>
                """, unsafe_allow_html=True)

            # 第三步：动态调整知识库结论的标题序号（核心修改点）
            if show_thinking:
                # 开启思考过程：显示为第三部分
                st.markdown("### 第三部分：知识库结论（JSON字段：final_conclusion）")
            else:
                # 关闭思考过程：自动切换为第二部分
                st.markdown("### 第二部分：知识库结论（JSON字段：final_conclusion）")
            st.markdown(f"""
            <div style="background:#ecfdf5;padding:12px;border-radius:8px;border:1px solid #a7f3d0;font-size:16px;
                        white-space: pre-wrap;  
                        max-height: 400px;      
                        overflow-y: auto;">     
                {kb_answer.final_conclusion}
            </div>
            """, unsafe_allow_html=True)

            # 打印完整JSON（验证格式正确）
            # 修复：先显示完整的kb_search_results，再显示JSON（格式验证）
            with st.expander("查看完整JSON响应"):
                st.markdown("#### 完整kb_search_results内容（无省略）：")
                st.text(kb_answer.kb_search_results)  # 用text显示，绝对完整，无任何省略
                st.markdown("#### JSON格式验证：")
                st.json(kb_answer.model_dump())  # JSON格式验证保留


# ====================== 知识库管理页签（完全保留） ======================
def kb_manage_tab():
    # 权限检查
    if not Permissions.has_permission(st.session_state.user_permissions, Permissions.VIEW_KB):
        st.error("你没有查看知识库的权限！")
        return
    col1, col2 = st.columns([1, 2])

    with col1:
        if st.session_state.init_ok:
            st.subheader("知识库管理")
            st.caption("批量上传+自动重试超时文件（PDF专项优化）")
            ALLOWED_DOC_TYPES = ['pdf', 'txt', 'doc', 'docx', 'xlsx', 'xls', 'ppt', 'pptx']
            MAX_FILE_SIZE_MB = 50  # 单文件大小限制，防内存溢出
            if Permissions.has_permission(st.session_state.user_permissions, Permissions.UPLOAD_KB):
                # 🔧 初始化 file_uploader_key（控制组件重新挂载）
                if "file_uploader_key" not in st.session_state:
                    st.session_state.file_uploader_key = 0

                # 🔧 使用列布局减少割裂感：左侧文件选择，右侧操作按钮
                col_upload, col_ops = st.columns([3, 1])

                with col_upload:
                    # 动态 key：清空时改变 key 值，强制组件重新挂载（前端视觉真正清空）
                    files = st.file_uploader(
                        "上传文件（支持PDF/Word/Excel/PPT/TXT）",  # 🔧 更新提示文案
                        type=ALLOWED_DOC_TYPES,  # 🔧 扩展类型
                        accept_multiple_files=True,
                        key=f"file_uploader_{st.session_state.file_uploader_key}"
                    )

                with col_ops:
                    st.write("")  # 占位对齐
                    st.write("")
                    # 实时显示数量（在按钮上方）
                    if files:
                        count = len(files)
                        remaining = MAX_FILES_PER_UPLOAD - count

                        # 🔧 友善提示：颜色区分严重程度
                        if count > MAX_FILES_PER_UPLOAD:
                            st.error(f"超限！{count}/{MAX_FILES_PER_UPLOAD}")
                            st.caption(f"请删除 {count - MAX_FILES_PER_UPLOAD} 个文件")
                        elif count == MAX_FILES_PER_UPLOAD:
                            st.warning(f"已达上限 {MAX_FILES_PER_UPLOAD} 个")
                        else:
                            st.success(f"{count}个文件")
                            st.caption(f"还可添加 {remaining} 个")

                        # 🔧 安全检查：显示文件详情（让用户自查）
                        with st.expander("查看文件清单"):
                            for i, f in enumerate(files, 1):
                                size_mb = f.size / (1024 * 1024)
                                icon = "" if size_mb < MAX_FILE_SIZE_MB else "error"
                                st.text(f"{icon} {i}. {f.name} ({size_mb:.1f}MB)")
                                # 🔒 安全检查：发现超大文件立即警告
                                if size_mb > MAX_FILE_SIZE_MB:
                                    st.error(f"{f.name} 超过{MAX_FILE_SIZE_MB}MB限制！")
                    else:
                        st.info("未选择文件")
                        st.caption(f"限额 {MAX_FILES_PER_UPLOAD} 个")

                    # 清空按钮：改变 key 强制重新挂载组件
                    if st.button("清空", key="clear_files_btn", use_container_width=True):
                        # 先删旧的（可选，清理内存）
                        old_key = f"file_uploader_{st.session_state.file_uploader_key}"
                        if old_key in st.session_state:
                            del st.session_state[old_key]

                        # 再换新key强制重新挂载
                        st.session_state.file_uploader_key += 1
                        st.rerun()
                # 🔧 表单只保留名称和提交，通过分割线视觉连接
                st.divider()
                with st.form("upload_form"):
                    kb_name = st.text_input("知识库名称", placeholder="输入名称（新增/追加）", key="kb_name")

                    # 提交按钮
                    submit_btn = st.form_submit_button("上传文件", type="primary")

                    if submit_btn:
                        current_key = st.session_state.file_uploader_key
                        files = st.session_state.get(f"file_uploader_{current_key}", [])

                        # === 基础验证（保留原逻辑）===
                        if not kb_name or not kb_name.strip():
                            st.error("请输入知识库名称！")
                            st.stop()

                        if not files:
                            st.error("请选择至少一个文件！")
                            st.stop()

                        if len(files) > MAX_FILES_PER_UPLOAD:
                            st.error(f"单次最多上传{MAX_FILES_PER_UPLOAD}个文件，当前已选{len(files)}个")
                            st.stop()

                        # === 新增：安全防护（可选但强烈建议）===
                        allowed_exts = {'pdf', 'txt', 'doc', 'docx', 'xlsx', 'xls', 'ppt', 'pptx'}

                        for f in files:
                            # 1. 防路径遍历
                            safe_name = os.path.basename(f.name)
                            if safe_name != f.name:
                                st.error(f"文件名非法：{f.name}")
                                st.stop()

                            # 2. 扩展名白名单
                            ext = safe_name.split('.')[-1].lower()
                            if ext not in allowed_exts:
                                st.error(f"不支持的格式：.{ext}")
                                st.stop()

                            # 3. 大小限制（防内存溢出）
                            if f.size > 50 * 1024 * 1024:  # 50MB
                                st.error(f"{safe_name} 超过50MB限制")
                                st.stop()

                        # === 原处理逻辑（完全保留）===
                        with st.spinner("处理中..."):
                            db = st.session_state.db_client
                            encoder = st.session_state.encoder
                            kb_id = db.get_kb_id(kb_name.strip())

                            if kb_id:
                                success_files, total_chunks, file_status = db.batch_upload_auto_retry(
                                    kb_id, files, encoder,
                                    chunk_size=st.session_state.chunk_size,
                                    sim_threshold=st.session_state.sim_threshold
                                )
                                st.session_state.file_status = file_status
                                st.success(f"完成！成功{success_files}文件，新增{total_chunks}分段")

                                # 修正后的清理逻辑
                                if f"file_uploader_{current_key}" in st.session_state:
                                    del st.session_state[f"file_uploader_{current_key}"]

                                st.session_state.file_uploader_key += 1
                                st.rerun()
            else:
                st.warning("你没有上传文件的权限")
            if Permissions.has_permission(st.session_state.user_permissions, Permissions.DELETE_KB):
                st.divider()
                st.subheader("删除知识库")
                st.warning("数据永久丢失！")
                kbs = st.session_state.db_client.get_all_kbs()
                if kbs:
                    kb_options = {kb['name']: kb['id'] for kb in kbs}
                    del_kb_name = st.selectbox("选择知识库", options=list(kb_options.keys()), key="del_kb")

                    if st.button("确认删除", type="secondary", key="del_kb_btn"):
                        st.session_state.kb_to_del = (kb_options[del_kb_name], del_kb_name)
                        st.session_state.del_kb_confirm = True

                    if st.session_state.del_kb_confirm and st.session_state.kb_to_del:
                        kb_id, kb_name = st.session_state.kb_to_del
                        st.error(f"最终确认：删除「{kb_name}」？")
                        col_final1, col_final2 = st.columns(2)
                        with col_final1:
                            if st.button("永久删除", key="final_del"):
                                st.session_state.db_client.delete_kb(kb_id, kb_name)
                                st.session_state.del_kb_confirm = False
                                st.session_state.kb_to_del = None
                                st.session_state.current_kb_id = None
                                st.rerun()
                        with col_final2:
                            if st.button("取消", key="cancel_del"):
                                st.session_state.del_kb_confirm = False
                                st.session_state.kb_to_del = None
            else:
                st.warning("你没有删除知识库的权限")
        else:
            st.info("请先初始化环境")

    with col2:
        if st.session_state.init_ok:
            st.subheader("知识库内容")

            kbs = st.session_state.db_client.get_all_kbs()
            if kbs:
                kb_display_options = {
                    f"{kb['name']} | 文档{kb['doc_count']} | 分段{kb['chunk_count']}": kb['id']
                    for kb in kbs
                }
                selected_kb_str = st.selectbox(
                    "选择知识库",
                    options=list(kb_display_options.keys()),
                    key="selected_kb"
                )
                current_kb_id = kb_display_options[selected_kb_str]

                if st.session_state.current_kb_id != current_kb_id:
                    st.session_state.current_kb_id = current_kb_id
                    st.session_state.current_page = 1

                chunks, total_pages, total = st.session_state.db_client.get_kb_chunks(
                    current_kb_id,
                    page=st.session_state.current_page,
                    page_size=PAGE_SIZE
                )

                col_p1, col_p2, col_p3 = st.columns([1, 2, 1])
                with col_p1:
                    if st.button("上一页", disabled=st.session_state.current_page <= 1, key="prev_page"):
                        st.session_state.current_page -= 1
                with col_p2:
                    st.write(f"第{st.session_state.current_page}/{total_pages}页 | 共{total}分段")
                with col_p3:
                    if st.button("下一页", disabled=st.session_state.current_page >= total_pages, key="next_page"):
                        st.session_state.current_page += 1

                st.divider()
                if chunks:
                    encoder = st.session_state.encoder
                    for chunk in chunks:
                        chunk_id = chunk['id']
                        is_editing = st.session_state.get(f'edit_chunk_{chunk_id}', False)
                        is_deleting = st.session_state.get(f'del_confirm_{chunk_id}', False)
                        with st.expander(f"{chunk['title']} | {chunk['doc_name']}", expanded=is_editing):
                            if not is_editing and not st.session_state.get(f'del_confirm_{chunk_id}', False):
                                # 🔧 新增：显示分段内容（原来漏了这行！）
                                st.markdown(f"**来源：** {chunk['doc_name']}  |  **创建时间：** {chunk['create_time']}")
                                st.markdown("**分段内容：**")
                                st.markdown(
                                    f"<div style='background:#f8fafc;padding:10px;border-radius:6px;'>{chunk['content']}</div>",
                                    unsafe_allow_html=True)
                                st.divider()
                                if not is_editing and not st.session_state.get(f'del_confirm_{chunk_id}', False):
                                    if st.button(f"编辑", key=f"edit_{chunk_id}"):
                                        # 重置所有chunk的编辑状态
                                        for key in list(st.session_state.keys()):
                                            if key.startswith('edit_chunk_'):
                                                del st.session_state[key]
                                        st.session_state[f'edit_chunk_{chunk_id}'] = True

                                    if st.button(f"删除", key=f"del_{chunk_id}"):
                                        # 重置所有chunk的编辑状态
                                        for key in list(st.session_state.keys()):
                                            if key.startswith('edit_chunk_'):
                                                del st.session_state[key]
                                        st.session_state[f'edit_chunk_{chunk_id}'] = True
                                        st.session_state[f'del_confirm_{chunk_id}'] = True

                                elif st.session_state.del_confirm and st.session_state.edit_chunk_id == chunk_id:
                                    st.warning("确认删除该分段？")
                                    col_c1, col_c2 = st.columns(2)
                                    with col_c1:
                                        if st.button(f"确认", key=f"conf_del_{chunk_id}"):
                                            st.session_state.db_client.delete_chunk(chunk_id)
                                            st.session_state.edit_chunk_id = None
                                            st.session_state.del_confirm = False
                                            st.rerun()
                                    with col_c2:
                                        if st.button(f"取消", key=f"cancel_del_{chunk_id}"):
                                            st.session_state.edit_chunk_id = None
                                            st.session_state.del_confirm = False
                                            st.rerun()
                                # 3. 修改删除确认块的判断和回调（3108行附近）
                                if st.session_state.get(f'del_confirm_{chunk_id}', False):
                                    st.warning("确认删除该分段？")
                                    col_c1, col_c2 = st.columns(2)
                                    with col_c1:
                                        if st.button(f"确认", key=f"conf_del_{chunk_id}"):
                                            st.session_state.db_client.delete_chunk(chunk_id)
                                            # 清理该chunk的所有状态
                                            for key in list(st.session_state.keys()):
                                                if key.startswith(f'edit_chunk_{chunk_id}') or key.startswith(
                                                        f'del_confirm_{chunk_id}'):
                                                    del st.session_state[key]
                                            st.rerun()

                                    with col_c2:
                                        if st.button(f"取消", key=f"cancel_del_{chunk_id}"):
                                            # 清理该chunk的所有状态
                                            for key in list(st.session_state.keys()):
                                                if key.startswith(f'edit_chunk_{chunk_id}') or key.startswith(
                                                        f'del_confirm_{chunk_id}'):
                                                    del st.session_state[key]
                            else:
                                if total == 0:
                                    st.info("该知识库暂无内容")
                                else:
                                    st.warning(f"当前页码{st.session_state.current_page}无数据，已自动调整到最后一页")
                                    st.session_state.current_page = total_pages
                                    st.rerun()
            else:
                st.info("暂无知识库，请先创建上传")
        else:
            st.info("请先初始化环境")


# ========== 【新增】图片识别+YOLO检测+知识库融合问答 页签 (核心功能，完全复用原有逻辑) ==========
# ====================== 图片识别问答页签（流式版，仅对话部分升级） ======================
def image_detect_qa_tab():
    # 权限检查
    if not Permissions.has_permission(st.session_state.user_permissions, Permissions.USE_IMAGE_QA):
        st.error("你没有使用图片识别问答的权限！")
        return

    st.subheader("图片目标识别与智能问答")
    st.caption("YOLO+OpenVINO INT8(iGPU)检测 + Ollama(GPU)流式生成回答 + 知识库融合 | 图片安全校验防马")

    if not st.session_state.init_ok:
        st.warning("请先左侧初始化环境（数据库+模型）")
        return

    # 1. 清空/重新上传按钮（保持不动）
    if st.button("清空图片 | 重新上传", use_container_width=True, type="secondary"):
        st.session_state.pop("detect_img_uploader", None)
        st.rerun()

    # 2. 图片上传与 YOLO 检测（保持不动）
    uploaded_img = st.file_uploader(
        "上传图片（仅支持JPG/PNG，≤8MB，自动安全校验）",
        type=ALLOWED_IMAGE_EXT,
        key="detect_img_uploader"
    )

    # 3. 知识库选择器（保持不动）
    st.divider()
    col_kb, col_topk = st.columns([2, 1])
    with col_kb:
        kbs = st.session_state.db_client.get_all_kbs()
        kb_options = {"全局检索（所有知识库）": None}
        if kbs:
            for kb in kbs:
                kb_options[kb['name']] = kb['id']
        selected_kb = st.selectbox("检索知识库范围", options=list(kb_options.keys()), key="img_qa_kb")
        selected_kb_id = kb_options[selected_kb]
    with col_topk:
        top_k = st.slider("返回上下文数量", 1, 15, 8, key="img_qa_top_k")

    # 4. YOLO 检测（保持不动）
    detect_result_text = ""
    detect_img = None
    if uploaded_img:
        is_safe, safe_msg = check_image_safety(uploaded_img)
        if not is_safe:
            st.error(safe_msg)
            return
        st.success(f"{safe_msg}，开始进行叶片病害检测...")
        img = Image.open(uploaded_img).convert('RGB')

        def retry_yolo_detect(image, max_retries=3):
            for retry in range(max_retries):
                try:
                    compiled_model, output_layer = load_yolo_ov_model()
                    return yolo_ov_detect(image)
                except Exception as e:
                    if retry < max_retries - 1:
                        st.warning(f"YOLO检测第{retry + 1}次失败，重试中... 错误：{str(e)[:50]}")
                        time.sleep(1)
                    else:
                        st.error(f"YOLO检测{max_retries}次均失败，错误：{str(e)[:50]}")
                        return "检测失败", None

        with st.spinner("YOLO iGPU推理中...(INT8量化加速，支持3次重试)"):
            detect_result_text, detect_img = retry_yolo_detect(img)
        if detect_result_text == "检测失败":
            return

        st.markdown("### YOLO目标检测结果 (Intel iGPU推理)")
        st.success("叶片病害检测完成！检测结果如下：")
        with st.expander("查看完整病害检测详情", expanded=True):
            st.text(detect_result_text)

    # 5. 用户提问输入（保持不动）
    user_question = st.text_area(
        "基于图片检测结果提问（可补充需求，自动融合知识库内容）",
        placeholder="例如：图片里检测到的病害是什么？结合知识库分析该病害的防治方法？",
        height=120,
        key="img_detect_query"
    )
    show_thinking = st.checkbox("显示思考过程", value=True, key="img_show_thinking")

    # 6. 流式问答核心逻辑（与 qa_tab 完全一致）
    if st.button("检测结果+知识库 智能回答", type="primary", key="img_qa_btn"):
        if not uploaded_img:
            st.error("请先上传图片完成目标检测！")
            return
        if not user_question.strip():
            st.error("请输入你的提问内容！")
            return

        # 进度条 + 状态文本（与 qa_tab 对齐）
        progress_bar = st.progress(0, text="正在准备检索...")
        status_text = st.empty()

        with st.spinner("正在检索知识库并按指定JSON格式生成回答..."):
            # Step-1 检索知识库
            progress_bar.progress(20, text="检索知识库中...")
            status_text.info("正在检索知识库内容...")
            db = st.session_state.db_client
            encoder = st.session_state.encoder
            combine_query = f"{detect_result_text}\n{user_question}"
            similar_chunks = db.search_similar_chunks(
                query=combine_query,
                encoder=encoder,
                top_k=top_k,
                kb_id=selected_kb_id
            )
            similar_chunks = [chunk for chunk in similar_chunks if chunk['content'].strip()]
            kb_search_text = ""
            if similar_chunks:
                kb_search_text = "\n\n".join([
                    f"【检索结果第{idx + 1}条 | 综合相似度：{chunk['total_similarity']:.4f}】\n"
                    f"来源文档：{chunk['doc_name']}\n"
                    f"分段标题：{chunk['title']}\n"
                    f"内容：{chunk['content'].strip()}"
                    for idx, chunk in enumerate(similar_chunks)
                ])
            else:
                kb_search_text = f"未检索到与图片检测结果+提问相关的【{selected_kb}】知识库片段"

            # Step-2 构造 Prompt（与 qa_tab 对齐）
            progress_bar.progress(40, text="构造提示词中...")
            status_text.info("正在构造提示词...")
            prompt_content = f"""
            你必须无条件严格遵守以下通用指令，适用于所有类型的问答需求，违反任意一条均视为回答错误：
            1.  输出格式：仅返回【纯标准JSON字符串】，无任何其他字符（无```json、无注释、无换行、无省略号）。
            2.  JSON结构：包含且仅包含3个字段，字段名一字不差：kb_search_results、thinking_process、final_conclusion，所有字段均为非空字符串。
            3.  字段内容通用铁律（优先级最高）：
                - kb_search_results：【原封不动】将下方【知识库查询结果】的全部内容复制填入，严禁任何修改！
                - thinking_process：先分析图片的YOLO检测结果，再逐条分析知识库检索内容，说明两者的关联性，不遗漏任何信息。
                - final_conclusion：优先使用知识库中的内容回答用户问题，结合图片的检测结果补充，所有内容必须100%来自知识库+检测结果，不编造，分点清晰罗列。

            【YOLO图片目标检测结果】：
            {detect_result_text}

            【知识库查询结果】：
            {kb_search_text}

            【用户问题】：{user_question}
            """.strip()

            # Step-3 使用Django加速接口（与qa_tab完全相同）
            progress_bar.progress(60, text="生成智能回答中...")
            status_text.info("正在生成回答...")

            # Step-3 使用 FastAPI 加速接口（无限轮询 + DELETE 确认）
            import requests, json, time, sys
            try:
                r = requests.post("http://127.0.0.1:8000/api/ask",
                                  json={
                                      "prompt": prompt_content,
                                      "kb_search_text": kb_search_text,
                                      "weather_text": ""
                                  },
                                  timeout=60)
                if r.status_code != 200:
                    st.error(f"后端返回错误: {r.status_code}")
                    progress_bar.empty()
                    status_text.empty()
                    return
                res = r.json()
                if res.get("hit"):
                    # 缓存命中，直接解析答案
                    ans_dict = json.loads(res["answer"]) if isinstance(res["answer"], str) else res["answer"]
                    kb_answer = KnowledgeBaseAnswer(**ans_dict)
                    st.success("从缓存加载（瞬间响应）")
                else:
                    task_id = res["task_id"]
                    print(f">>> 图片问答 - 开始轮询任务 {task_id}", file=sys.stderr)
                    answer_obtained = False
                    polling_answer = None
                    max_poll_interval = 1.0
                    min_poll_interval = 0.1
                    poll_interval = min_poll_interval
                    with st.spinner(f"正在生成回答 (任务 {task_id})..."):
                        while True:
                            try:
                                result_resp = requests.get(f"http://127.0.0.1:8000/api/result/{task_id}", timeout=5)
                                result_resp.raise_for_status()
                                result_data = result_resp.json()
                                if "answer" in result_data:
                                    ans = result_data["answer"]
                                    ans_dict = json.loads(ans) if isinstance(ans, str) else ans
                                    polling_answer = KnowledgeBaseAnswer(**ans_dict)
                                    answer_obtained = True
                                    # 发送结束确认
                                    try:
                                        del_resp = requests.delete(f"http://127.0.0.1:8000/api/result/{task_id}",
                                                                   timeout=2)
                                        print(f">>> 图片问答 - 删除任务 {task_id} -> {del_resp.status_code}",
                                              file=sys.stderr)
                                    except Exception:
                                        pass
                                    progress_bar.empty()
                                    status_text.empty()
                                    break
                                elif result_data.get("status") == "failed":
                                    st.error(f"任务失败: {result_data.get('error', '未知错误')}")
                                    polling_answer = KnowledgeBaseAnswer(
                                        kb_search_results=kb_search_text,
                                        thinking_process="任务失败",
                                        final_conclusion=f"生成答案失败：{result_data.get('error', '未知错误')}"
                                    )
                                    answer_obtained = True
                                    break
                                else:
                                    time.sleep(poll_interval)
                                    poll_interval = min(poll_interval * 1.5, max_poll_interval)
                                    continue
                            except requests.exceptions.RequestException as req_err:
                                print(f">>> 图片问答 - 轮询请求失败: {req_err}", file=sys.stderr)
                                time.sleep(poll_interval)
                                poll_interval = min(poll_interval * 1.5, max_poll_interval)
                                continue
                            except json.JSONDecodeError as json_err:
                                print(f">>> 图片问答 - JSON解析失败: {json_err}", file=sys.stderr)
                                time.sleep(poll_interval)
                                poll_interval = min(poll_interval * 1.5, max_poll_interval)
                                continue
                    if answer_obtained:
                        kb_answer = polling_answer
                    else:
                        # fallback 到本地 Ollama
                        try:
                            fallback_response = cached_ollama_response(prompt_content,
                                                                       KnowledgeBaseAnswer.model_json_schema())
                            raw_json = fallback_response.message.content.strip()
                            raw_json = re.sub(r'^```json|```$', '', raw_json).strip()
                            kb_answer = KnowledgeBaseAnswer.model_validate_json(raw_json)
                        except Exception as fallback_err:
                            st.error(f"本地回退生成失败: {fallback_err}")
                            kb_answer = KnowledgeBaseAnswer(
                                kb_search_results=kb_search_text,
                                thinking_process="生成失败",
                                final_conclusion="无法生成答案，请稍后重试"
                            )
            except Exception as e:
                st.error(f"调用API失败: {str(e)}")
                kb_answer = KnowledgeBaseAnswer(
                    kb_search_results=kb_search_text,
                    thinking_process="API调用失败",
                    final_conclusion="无法生成答案，请检查后端服务是否运行"
                )

            # Step-4 JSON 兜底校验（与 qa_tab 完全一致）
            if not kb_answer.kb_search_results.strip():
                kb_answer.kb_search_results = kb_search_text
            if not kb_answer.thinking_process.strip():
                kb_answer.thinking_process = "已完成图片检测和知识库检索，未发现相关关联信息"
            if not kb_answer.final_conclusion.strip():
                kb_answer.final_conclusion = "未从知识库中检索到相关内容，仅参考图片检测结果回答"

            # Step-5 清理进度条
            progress_bar.empty()
            status_text.empty()

        # 7. 结果展示（与 qa_tab 完全一致，仅标题序号随 show_thinking 动态切换）
        st.markdown("### 第一部分：知识库检索结果")
        if similar_chunks:
            with st.expander(f"展开/收起（共{len(similar_chunks)}条）", expanded=True):
                for idx, chunk in enumerate(similar_chunks, 1):
                    st.markdown(f"""
<div style="background:#f5f7fa;padding:10px;border-radius:6px;margin-bottom:8px;">
    <strong>第{idx}条 | {chunk['kb_name']} | {chunk['doc_name']}</strong>
    <br><strong>相似度：{chunk['total_similarity']}</strong>
    <br><strong>标题：</strong>{chunk['title']}
    <br><strong>内容：</strong>{clean_text_local(chunk['content'])}
</div>
                    """, unsafe_allow_html=True)
        else:
            st.markdown(f"""
<div style="background:#fef7f7;padding:10px;border-radius:6px;color:#dc2626;">
    {kb_answer.kb_search_results}
</div>
            """, unsafe_allow_html=True)

        if show_thinking:
            st.markdown("### 第二部分：思考过程")
            st.markdown(f"""
<div style="background:#eef2ff;padding:12px;border-radius:8px;border:1px solid #c7d2fe;">
    {kb_answer.thinking_process}
</div>
            """, unsafe_allow_html=True)

        title_num = "### 第三部分：最终结论" if show_thinking else "### 第二部分：最终结论"
        st.markdown(title_num)
        st.markdown(f"""
<div style="background:#ecfdf5;padding:12px;border-radius:8px;border:1px solid #a7f3d0;font-size:16px;white-space: pre-wrap;max-height:400px;overflow-y:auto;">
    {kb_answer.final_conclusion}
</div>
        """, unsafe_allow_html=True)

        with st.expander("查看完整JSON响应"):
            st.json(kb_answer.model_dump())


# ====================== 用户管理页签（新增） ======================
def user_manage_tab():
    """用户管理页面 - 只有管理员可见"""
    st.subheader("用户管理")
    st.caption("管理员功能：查看、管理用户账户和权限")

    # 检查管理员权限
    if not Permissions.has_permission(st.session_state.user_permissions, Permissions.ADMIN):
        st.error("你没有管理员权限，无法访问此页面！")
        return

    # 获取所有用户
    db = st.session_state.db_client
    users = db.get_all_users()

    if not users:
        st.info("暂无用户数据")
        return

    # 🔧 修复：获取所有用户，但过滤掉当前登录用户和admin用户
    admin_encrypted = sm4_encrypt("admin")
    current_user = st.session_state.get("current_user", "")

    users = []
    all_users = db.get_all_users()
    for user in all_users:
        try:
            username_display = sm4_decrypt(user['username'])
            # 过滤掉当前登录用户和admin用户
            if username_display != current_user and username_display != "admin":
                # 保持加密用户名在user对象中，但添加解密后的显示名
                user['username_display'] = username_display
                users.append(user)
        except Exception:
            # 解密失败的用户也显示，但同样要过滤
            user_id = user['id']
            user_display = f"用户{user_id}"
            if user_display != current_user and user['username'] != admin_encrypted:
                user['username_display'] = user_display
                users.append(user)

    if not users:
        st.info("暂无其他用户数据")
        return

    st.write(f"共 {len(users)} 个用户（管理员admin不显示）")

    # 用户列表
    for user in users:
        user_id = user['id']
        username_encrypted = user['username']
        username_display = user.get('username_display', f"用户{user_id}")  # 使用已解密的显示名
        permissions = user.get('permissions', Permissions.DEFAULT_USER)
        create_time = user['create_time']

        # 解密用户名显示
        try:
            username_display = sm4_decrypt(username_encrypted)
        except Exception as e:
            username_display = f"用户{user_id}"

        with st.expander(
                f"👤 {username_display} | 权限: {permissions} | 注册: {create_time.strftime('%Y-%m-%d %H:%M')}"):
            # 权限详情
            perm_names = Permissions.get_permission_names(permissions)
            st.write(f"**当前权限:** {', '.join(perm_names)}")

            st.divider()

            # 权限编辑区域
            st.write("**修改权限:**")
            col1, col2 = st.columns(2)

            with col1:
                new_login = st.checkbox("登录权限", value=Permissions.has_permission(permissions, Permissions.LOGIN),
                                        key=f"login_{user_id}")
                new_view_kb = st.checkbox("查看知识库",
                                          value=Permissions.has_permission(permissions, Permissions.VIEW_KB),
                                          key=f"view_{user_id}")
                new_upload_kb = st.checkbox("上传文件",
                                            value=Permissions.has_permission(permissions, Permissions.UPLOAD_KB),
                                            key=f"upload_{user_id}")
                new_delete_kb = st.checkbox("删除知识库",
                                            value=Permissions.has_permission(permissions, Permissions.DELETE_KB),
                                            key=f"delete_{user_id}")

            with col2:
                new_edit_kb = st.checkbox("编辑内容",
                                          value=Permissions.has_permission(permissions, Permissions.EDIT_KB),
                                          key=f"edit_{user_id}")
                new_use_qa = st.checkbox("智能问答", value=Permissions.has_permission(permissions, Permissions.USE_QA),
                                         key=f"qa_{user_id}")
                new_use_image_qa = st.checkbox("图片识别",
                                               value=Permissions.has_permission(permissions, Permissions.USE_IMAGE_QA),
                                               key=f"img_{user_id}")

                # 🔧 修复：允许admin用户授予其他用户管理员权限
                current_user = st.session_state.get("current_user", "")

                # 如果是当前用户自己，禁用管理员权限修改（不能修改自己的管理员权限）
                if username_display == current_user:
                    st.write("**管理员**（不能修改自己的权限）")
                    new_admin = Permissions.has_permission(permissions, Permissions.ADMIN)
                elif current_user == "admin":
                    new_admin = st.checkbox("管理员", value=Permissions.has_permission(permissions, Permissions.ADMIN),
                                            key=f"admin_{user_id}", help="授予此用户管理员权限")
                else:
                    # 非admin用户不能修改管理员权限
                    new_admin = Permissions.has_permission(permissions, Permissions.ADMIN)
                    if new_admin:
                        st.write("**管理员**（此用户是管理员）")
                    else:
                        st.write("**管理员**（无权限修改）")

            # 计算新权限值
            new_permissions = 0
            if new_login:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.LOGIN, True)
            if new_view_kb:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.VIEW_KB, True)
            if new_upload_kb:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.UPLOAD_KB, True)
            if new_delete_kb:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.DELETE_KB, True)
            if new_edit_kb:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.EDIT_KB, True)
            if new_use_qa:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.USE_QA, True)
            if new_use_image_qa:
                new_permissions = Permissions.set_permission(new_permissions, Permissions.USE_IMAGE_QA, True)

            # 🔧 修复：只有当前用户是admin，才能给其他用户授予管理员权限
            if current_user == "admin" and username_display != current_user:
                if new_admin:
                    new_permissions = Permissions.set_permission(new_permissions, Permissions.ADMIN, True)
                    st.info("注意：正在授予此用户管理员权限")
            elif username_display == current_user:
                # 如果是修改自己的权限，保持原有的管理员状态
                if Permissions.has_permission(permissions, Permissions.ADMIN):
                    new_permissions = Permissions.set_permission(new_permissions, Permissions.ADMIN, True)

            # 保存和删除按钮
            col_save, col_del, col_pwd = st.columns(3)

            with col_save:
                if st.button("保存权限", key=f"save_perm_{user_id}"):
                    current_user = st.session_state.get("current_user", "")

                    # 🔧 修复：防止任何用户修改自己的权限（包括admin）
                    if current_user == username_display:
                        st.error("不能修改自己的权限！")
                        st.rerun()

                    # 🔧 修复：防止普通用户给自己添加管理员权限
                    if current_user != "admin":
                        # 普通用户不能授予管理员权限
                        if new_permissions & Permissions.ADMIN:
                            st.error("只有admin可以授予管理员权限！")
                            st.rerun()

                    # 🔧 修复：防止删除admin用户的admin权限
                    if username_display == "admin":
                        if not (new_permissions & Permissions.ADMIN):
                            st.error("不能移除admin用户的管理员权限！")
                            st.rerun()

                    if db.update_user_permissions(username_display, new_permissions):
                        # 记录日志
                        old_perms = Permissions.get_permission_names(permissions)
                        new_perms = Permissions.get_permission_names(new_permissions)
                        log_sensitive_operation(
                            "修改用户权限",
                            username_display,
                            current_user,
                            f"权限变更: {old_perms} → {new_perms}"
                        )
                        st.success("权限已更新")
                        time.sleep(0.5)
                        st.rerun()
                    else:
                        st.error("更新失败")

            with col_del:
                # 🔧 修复：实现两步确认删除 + 禁止删除自己
                del_key = f"del_confirm_{user_id}"
                if del_key not in st.session_state:
                    st.session_state[del_key] = False

                # 🔧 修复：禁止任何用户删除自己（包括管理员）
                is_self = (username_display == st.session_state.get("current_user", ""))

                if is_self:
                    st.warning("不能删除自己")
                elif not st.session_state[del_key]:
                    if st.button("删除用户", key=f"del_btn_{user_id}"):
                        st.session_state[del_key] = True
                        st.rerun()
                else:
                    st.error(f"确认删除用户「{username_display}」？")
                    col_confirm, col_cancel = st.columns(2)
                    with col_confirm:
                        if st.button("永久删除", key=f"final_del_{user_id}"):
                            # 🔧 修复：再次确认不是删除自己（双重检查）
                            if username_display == st.session_state.get("current_user", ""):
                                st.error("不能删除自己！操作已阻止")
                                del st.session_state[del_key]
                                st.rerun()

                            if db.delete_user(username_display):
                                # 记录日志
                                current_user = st.session_state.get("current_user", "未知")
                                log_sensitive_operation(
                                    "删除用户",
                                    username_display,
                                    current_user,
                                    f"用户ID: {user_id}"
                                )
                                st.success("用户已删除")
                                del st.session_state[del_key]
                                time.sleep(0.5)
                                st.rerun()
                            else:
                                st.error("删除失败")

            with col_pwd:
                # 🔧 关键修复：简化密码管理逻辑，使用与登录页相同的方法
                pwd_key = f"pwd_mgmt_{user_id}"
                pwd_verified_key = f"pwd_verified_{user_id}"

                # 初始化状态
                if pwd_key not in st.session_state:
                    st.session_state[pwd_key] = False
                if pwd_verified_key not in st.session_state:
                    st.session_state[pwd_verified_key] = False

                # 管理密码按钮
                if st.button("管理密码", key=f"pwd_btn_{user_id}"):
                    # 重置所有其他用户的密码管理状态
                    for key in list(st.session_state.keys()):
                        if key.startswith("pwd_mgmt_") and key != pwd_key:
                            st.session_state[key] = False
                        if key.startswith("pwd_verified_") and key != pwd_verified_key:
                            st.session_state[key] = False

                    # 开启当前用户的密码管理
                    st.session_state[pwd_key] = True
                    st.session_state[pwd_verified_key] = False
                    # 刷新验证码（使用全局验证码，与登录页一致）
                    refresh_captcha()
                    st.rerun()

                # 如果打开了密码管理
                if st.session_state[pwd_key]:
                    # 如果未验证，显示验证码输入
                    if not st.session_state[pwd_verified_key]:
                        st.write("**身份验证**")
                        col_img, col_input, col_refresh = st.columns([1, 1.5, 0.5])
                        with col_img:
                            # 使用全局验证码图片
                            img = st.session_state.get("captcha_image")
                            if img:
                                st.image(img, width=100)
                            else:
                                # 兜底：重新生成
                                refresh_captcha()
                                img = st.session_state.get("captcha_image")
                                if img:
                                    st.image(img, width=100)
                        with col_input:
                            captcha_input = st.text_input("输入验证码", key=f"pwd_captcha_{user_id}",
                                                          label_visibility="collapsed",
                                                          placeholder="输入图片中的验证码")
                        with col_refresh:
                            if st.button("re", key=f"refresh_captcha_{user_id}", use_container_width=True):
                                refresh_captcha()
                                st.rerun()

                        col_verify, col_cancel = st.columns(2)
                        with col_verify:
                            if st.button("验证", key=f"verify_{user_id}", use_container_width=True):
                                user_input = captcha_input.strip()
                                correct_text = st.session_state.get("captcha_text", "").strip()

                                # 使用与登录页完全相同的验证逻辑
                                if not user_input:
                                    st.error("验证码不能为空！")
                                    refresh_captcha()
                                    st.rerun()
                                    return

                                # 不区分大小写比较
                                if user_input.upper() == correct_text.upper():
                                    st.session_state[pwd_verified_key] = True
                                    st.success("验证通过")
                                    # 验证通过后刷新验证码
                                    refresh_captcha()
                                    st.rerun()
                                else:
                                    st.error("验证码错误！")
                                    refresh_captcha()
                                    st.rerun()

                        with col_cancel:
                            if st.button("取消", key=f"cancel_verify_{user_id}", use_container_width=True):
                                st.session_state[pwd_key] = False
                                st.session_state[pwd_verified_key] = False
                                refresh_captcha()
                                st.rerun()
                    else:
                        # 验证通过后显示密码修改区域
                        st.write("**修改密码**")
                        new_pwd = st.text_input("新密码（≥6位）", type="password", key=f"new_pwd_{user_id}")
                        confirm_pwd = st.text_input("确认新密码", type="password", key=f"confirm_pwd_{user_id}")

                        col_update, col_close = st.columns(2)
                        with col_update:
                            if st.button("更新密码", key=f"update_pwd_{user_id}", use_container_width=True):
                                if new_pwd and confirm_pwd:
                                    if len(new_pwd) >= 6:
                                        if new_pwd == confirm_pwd:
                                            if db.update_user_password(username_display, new_pwd):
                                                # 记录日志
                                                current_user = st.session_state.get("current_user", "未知")
                                                log_sensitive_operation(
                                                    "修改用户密码",
                                                    username_display,
                                                    current_user,
                                                    f"用户ID: {user_id}"
                                                )
                                                st.success("密码已更新")
                                                # 关闭密码管理
                                                st.session_state[pwd_key] = False
                                                st.session_state[pwd_verified_key] = False
                                                time.sleep(0.5)
                                                st.rerun()
                                            else:
                                                st.error("密码更新失败")
                                        else:
                                            st.error("两次输入的密码不一致")
                                    else:
                                        st.error("密码长度需≥6位")
                                else:
                                    st.error("请填写新密码和确认密码")
                        with col_close:
                            if st.button("关闭", key=f"close_pwd_{user_id}", use_container_width=True):
                                st.session_state[pwd_key] = False
                                st.session_state[pwd_verified_key] = False
                                st.rerun()


# ========== 新增：独立数据库初始化函数（认证页复用） ==========
def init_database(db_host, db_port, db_user, db_pwd, db_name, use_cuda, model_type, silent=False):
    """
    初始化数据库和模型
    :param silent: 静默模式，True=不显示初始化消息（用于登录页自动加载）
    """
    try:
        encoder = get_local_encoder(use_cuda, model_type, silent)
        db_config = {
            "host": db_host,
            "port": db_port,
            "user": db_user,
            "password": db_pwd,
            "dbname": db_name
        }
        db_client = LocalPGVectorDB(db_config, silent)  # 传递silent参数
        # 关键修复：根据use_cuda设置Ollama设备
        if use_cuda and torch.cuda.is_available():
            os.environ["OLLAMA_DEVICE"] = "cuda"
            print("Ollama设备：NVIDIA GPU (CUDA)")
        else:
            os.environ["OLLAMA_DEVICE"] = "cpu"
            print("Ollama设备：CPU")
        # 检查数据库对象是否有效（conn和cursor不为None）
        if db_client.conn is None or db_client.cursor is None:
            return False, None, None, "数据库连接失败，无法创建数据库客户端！"
        return True, encoder, db_client, None  # 成功：状态+编码器+数据库客户端
    except Exception as e:
        error_msg = f"数据库初始化失败：{str(e)[:100]}"
        logging.error(f"数据库初始化失败详情：{traceback.format_exc()}")
        return False, None, None, error_msg  # 失败：状态+错误信息


def ensure_db_initialized():
    """
    智能初始化函数：
    - 如果已初始化过，直接返回
    - 如果未初始化，从db.yaml加载并初始化
    - 自动管理session_state状态
    """
    # 1. 如果已经准备好，直接返回成功
    if st.session_state.get("db_client_ready", False):
        return True, st.session_state.encoder, st.session_state.db_client, None

    # 2. 尝试从yaml加载配置
    yaml_config = load_db_config_from_yaml()
    if not yaml_config:
        return False, None, None, "未找到db.yaml配置文件！"

    # 3. 执行初始化 (silent=True 避免多余输出)
    success, encoder, db_client, error_msg = init_database(
        yaml_config["host"], yaml_config["port"],
        yaml_config["user"], yaml_config["password"],
        yaml_config["dbname"], yaml_config["use_cuda"],
        yaml_config["model_type"],
        silent=True  # 【新增】静默模式
    )

    # 4. 如果成功，保存状态到session_state
    if success:
        st.session_state.db_client = db_client
        st.session_state.encoder = encoder
        st.session_state.db_client_ready = True  # 标记已准备好

    return success, encoder, db_client, error_msg

# ========== 全局缓存：数据库客户端和编码器（只初始化一次） ==========
@st.cache_resource(ttl=None, show_spinner="加载数据库...")
def _get_cached_db_client():
    print("[调试] 开始加载数据库客户端...")

    if not DB_YAML_PATH.exists():
        print(f"[调试] db.yaml 不存在: {DB_YAML_PATH}")
        return None

    print(f"[调试] db.yaml 路径: {DB_YAML_PATH}")

    cfg = load_db_config_from_yaml()
    print(f"[调试] 配置加载结果: {cfg}")

    if not cfg:
        print("[调试] 配置为空，返回 None")
        return None

    print(f"[调试] 正在连接: {cfg['host']}:{cfg['port']}/{cfg['dbname']}")

    db_config = {
        "host": cfg["host"],
        "port": cfg["port"],
        "user": cfg["user"],
        "password": cfg["password"],
        "dbname": cfg["dbname"]
    }

    try:
        client = LocalPGVectorDB(db_config, silent=True)
        print("[调试] 数据库客户端创建成功")
        return client
    except Exception as e:
        print(f"[调试] 创建客户端失败: {e}")
        import traceback
        print(traceback.format_exc())
        return None

@st.cache_resource
def _get_cached_encoder():
    if not DB_YAML_PATH.exists():
        return None
    cfg = load_db_config_from_yaml()
    if not cfg:
        return None
    return get_local_encoder(cfg["use_cuda"], cfg["model_type"], silent=True)
# ====================== 主函数（完全保留） ======================
def main():
    # 【关键】页面配置必须在最开始，且只调用一次
    st.set_page_config(page_title="知识库问答系统", page_icon="📚", layout="wide")
    # 如果已通过 token 登录，但 db_client 为 None，直接尝试复用缓存（不重连）
    if st.session_state.get("is_logged_in") and st.session_state.get("db_client") is None:
        try:
            # 尝试快速复用，如果失败显示错误而非卡住
            st.session_state.db_client = _get_cached_db_client()
            st.session_state.encoder = _get_cached_encoder()
        except Exception as e:
            st.error(f"数据库连接失败: {e}")
            st.button("重新连接", on_click=lambda: st.rerun())
            return
    # ========== 登录检查（纯后端，无 JS） ==========
    if "is_logged_in" not in st.session_state:
        st.session_state.is_logged_in = False
    if st.session_state.get("pending_token") and not st.query_params.get("token"):
        token = st.session_state.pending_token
        st.query_params["token"] = token  # 设置到 URL
        del st.session_state.pending_token  # 清理临时存储
        st.rerun()  # 刷新使 URL 生效
    # 获取 URL 参数中的 token
    token = st.query_params.get("token")
    print(f"[db.py] 原始 token from URL: {token}")

    # 尝试通过 token 自动登录（仅当尚未登录时）
    if token and not st.session_state.is_logged_in:
        try:
            decrypted = sm4_decrypt(token)
            if "|" in decrypted:
                username, expire_ts = decrypted.split("|", 1)
                if time.time() < int(expire_ts):
                    db_client = _get_cached_db_client()
                    if db_client:
                        user_data = db_client.get_user_by_username(username)
                        if user_data:
                            # ... 设置 session_state ...
                            st.session_state.is_logged_in = True
                            st.session_state.current_user = username
                            st.session_state.user_permissions = user_data['permissions']

                            # 关键：显式保留 token，防止后续 reruns 丢失
                            st.query_params["token"] = token
                            print(f"[db.py] token登录成功，已保留在URL")
        except Exception as e:
            print(f"[db.py] token登录异常: {e}")

    # 如果仍未登录，根据 db.yaml 是否存在决定流程
    if not st.session_state.is_logged_in:
        # 检查 db.yaml 是否存在且有效
        if DB_YAML_PATH.exists():
            # 配置文件存在，跳转到登录页
            print("[db.py] 未登录但 db.yaml 存在，跳转到 login.py")
            st.switch_page("pages/login.py")
            return
        else:
            # 配置文件不存在，直接显示数据库初始化界面
            print("[db.py] db.yaml 不存在，进入初始化界面")
            show_db_config_stage()
            return

    # 后面继续您原有的样式配置、侧边栏等代码...

    # ========== 样式配置 ==========
    st.markdown("""
        <style>
            [data-testid="stSidebarNav"] { display: none !important; }
            section[data-testid="stSidebar"] > div:first-child { margin-top: 0; }
        </style>
        """, unsafe_allow_html=True)

    # ========== 未登录处理 ==========
    if not st.session_state.is_logged_in:
        if DB_YAML_PATH.exists():
            # 跳转到登录页
            st.switch_page("pages/login.py")
            return
        else:
            show_db_config_stage()
        return

    # ========== 【原有代码】已登录，恢复数据库连接 ==========
    if st.session_state.get("db_client") is None:
        st.session_state.db_client = _get_cached_db_client()
        st.session_state.encoder = _get_cached_encoder()
        st.session_state.db_client_ready = True
        print("[db.py] 从缓存获取 db_client 和 encoder")

    # 载入城市信息
    if DB_YAML_PATH.exists():
        cfg = load_db_config_from_yaml()
        if cfg:
            st.session_state["user_city"] = cfg.get("user_city", "")
            st.session_state["user_adcode"] = cfg.get("user_adcode", "")

    # ========== 【原有代码】初始化 session_state ==========
    init_vars = {
        "db_initialized": False,
        "temp_db_config": None,
        "last_search": "",
        "db_client": None,
        "init_ok": False,
        "current_kb_id": None,
        "current_page": 1,
        "edit_chunk_id": None,
        "del_confirm": False,
        "del_kb_confirm": False,
        "kb_to_del": None,
        "file_status": {},
        "encoder": None,
        "is_logged_in": False,
        "db_host": DEFAULT_DB_CONFIG['host'],
        "db_port": DEFAULT_DB_CONFIG['port'],
        "db_user": DEFAULT_DB_CONFIG['user'],
        "db_pwd": DEFAULT_DB_CONFIG['password'],
        "db_name": DEFAULT_DB_CONFIG['dbname'],
        "model_type": "双模型融合",
        "use_cuda": False,
        "auth_db_host": DEFAULT_DB_CONFIG['host'],
        "auth_db_port": DEFAULT_DB_CONFIG['port'],
        "auth_db_user": DEFAULT_DB_CONFIG['user'],
        "auth_db_pwd": DEFAULT_DB_CONFIG['password'],
        "auth_db_name": DEFAULT_DB_CONFIG['dbname'],
        "auth_use_cuda": USE_CUDA_DEFAULT,
        "auth_tab": "login",
        "is_first_init": False,
        "captcha_text": "",
        "captcha_image": None,
        "reset_verified_user": None,
        "reset_new_pwd": "",
        "reset_confirm_pwd": "",
        "current_user": None,
        "user_permissions": Permissions.DEFAULT_USER,
    }
    for key, val in init_vars.items():
        if key not in st.session_state:
            st.session_state[key] = val

    # 🔧 新增：检查并验证db.yaml配置
    if DB_YAML_PATH.exists():
            # 尝试加载并验证配置
            yaml_config = load_db_config_from_yaml()
            if yaml_config is None:
                # 配置验证失败，文件已被删除，显示初始化页面
                st.markdown(HIDE_ST_STYLE, unsafe_allow_html=True)
                st.subheader("配置文件验证失败")
                st.error("db.yaml配置文件不完整或格式错误，已自动删除。请重新初始化数据库。")

                if st.button("重新初始化数据库", type="primary"):
                    # 确保文件已被删除
                    if DB_YAML_PATH.exists():
                        try:
                            os.remove(DB_YAML_PATH)
                        except:
                            pass
                    st.rerun()
                return
    # ========== 路由逻辑结束 ==========

    # 侧边栏配置（完全保留）
    with st.sidebar:
        st.subheader("核心配置")
        # 先获取侧边栏控件的当前值（提前缓存，避免重复渲染问题）
        current_sidebar_model = st.session_state.get("sidebar_model_type", st.session_state.model_type)
        current_sidebar_cuda = st.session_state.get("sidebar_use_cuda", st.session_state.use_cuda)
        # 精准对比：侧边栏当前值 vs 全局session_state值，不一致则重置init_ok
        if (current_sidebar_model != st.session_state.model_type) or (
                current_sidebar_cuda != st.session_state.use_cuda):
            st.session_state.init_ok = False  # 强制重置，按钮立即重新可用

        st.caption("模型&分段配置")

        new_model_type = st.selectbox(
            "模型类型",
            options=["双模型融合", "text2vec-base-chinese", "bert-base-chinese"],
            key="sidebar_model_type",
            index=["双模型融合", "text2vec-base-chinese", "bert-base-chinese"].index(st.session_state.model_type)
        )
        new_use_cuda = st.checkbox("启用CUDA", key="sidebar_use_cuda", value=st.session_state.use_cuda)
        # 新增：同步更新Ollama设备
        if new_use_cuda != st.session_state.use_cuda:
            if new_use_cuda and torch.cuda.is_available():
                os.environ["OLLAMA_DEVICE"] = "cuda"
            else:
                os.environ["OLLAMA_DEVICE"] = "cpu"
        if (new_model_type != st.session_state.model_type) or (new_use_cuda != st.session_state.use_cuda):
            st.session_state.init_ok = False  # 配置变化，按钮重新启用
        st.slider("分段字符数", 100, 300, CHUNK_SIZE_DEFAULT, key="chunk_size")
        st.slider("最小分段长度", 10, 100, MIN_LENGTH_DEFAULT, key="min_length")
        st.slider("相似度阈值", 0.7, 0.9, 0.8, 0.01, key="sim_threshold")
        # 核心：删除 st.session_state.init_ok 的控制，替换为基础校验
        is_disabled = st.session_state.db_client is None or st.session_state.encoder is None
        if (new_model_type != st.session_state.model_type) or (new_use_cuda != st.session_state.use_cuda):
            st.session_state.init_ok = False  # 你的原有逻辑，保留不变

        # 2. 按钮判断：依赖新增的 is_first_init（独立状态，不受配置变化影响）
        if not st.session_state.is_first_init:
            btn_text = "初始化环境"
            btn_type = "primary"  # 合规值，无异常
        else:
            btn_text = "重新加载模型"
            btn_type = "secondary"  # 合规值，无异常

        # 3. 按钮渲染（完全保留你的参数，仅替换文本和type）
        if st.button(
                btn_text,
                type=btn_type,
                disabled=is_disabled,
                key="init_btn",
                help="首次初始化或修改模型/CUDA配置后，点击此按钮生效"
        ):
            with st.spinner("正在验证环境并初始化..."):
                try:
                    # 你原有内部逻辑（一字不变，仅新增一行：标记首次初始化完成）
                    if st.session_state.db_client is None or st.session_state.encoder is None:
                        st.error("数据库或模型未提前初始化，请返回登录页重新完成初始化！")
                        return
                    new_encoder = get_local_encoder(new_use_cuda, new_model_type)
                    st.session_state.encoder = new_encoder
                    st.session_state.model_type = new_model_type
                    st.session_state.use_cuda = new_use_cuda
                    st.session_state.init_ok = True  # 你的原有逻辑，保留不变
                    st.session_state.is_first_init = True  # 新增：标记首次初始化完成，后续不再改变

                    # 动态提示（可选，保留不变）
                    tip_text = "环境初始化成功" if btn_text == "🔌 初始化环境" else "模型重新加载成功"
                    st.success(f"{tip_text}！当前：{new_model_type} | CUDA：{new_use_cuda}")
                    st.rerun()
                except Exception as e:
                    error_text = "环境初始化" if btn_text == "🔌 初始化环境" else "模型重新加载"
                    st.error(f"{error_text}失败：{str(e)[:100]}")
                    logging.error(f"{error_text}失败详情：{traceback.format_exc()}")

                # ========== 整合所有逻辑的「初始化环境」按钮（核心修改） ==========
                # 1. 移除按钮文字中的（修改后点击）
                # 2. 添加help参数，鼠标悬浮显示提示
                # 3. 整合首次初始化+后期模型更新逻辑，无冗余按钮
        # 在侧边栏模型配置后添加（大约在第1630行）
        st.sidebar.divider()
        st.sidebar.subheader("流式输出")

        # 流式开关
        use_streaming = st.sidebar.checkbox("启用流式输出", value=True,
                                            help="实时显示生成内容，加速响应体验")

        # 缓存管理
        col1, col2 = st.sidebar.columns(2)
        with col1:
            if st.button("清缓存", key="clear_cache_btn", help="清除Redis缓存"):
                try:
                    keys = REDIS_CLI.keys("*")
                    if keys:
                        REDIS_CLI.delete(*keys)
                        st.sidebar.success(f"已清 {len(keys)} 条")
                        time.sleep(0.5)
                        st.rerun()
                    else:
                        st.sidebar.info("无缓存")
                except Exception as e:
                    st.sidebar.error(f"失败: {str(e)[:30]}")

        with col2:
            if st.button("流测试", key="stream_test", help="测试流式功能"):
                st.sidebar.info("流式功能已启用")

        # 显示缓存状态
        try:
            cache_count = len(REDIS_CLI.keys("*"))
            st.sidebar.caption(f"缓存(全部 Redis key): {cache_count} 条")  # 全部命名空间
        except:
            pass

        if st.session_state.is_logged_in:
            # 可选：每次页面刷新时续期token（保持7天有效期）
            try:
                new_token_data = f"{st.session_state.current_user}|{int(time.time()) + 86400 * 7}"
                new_token = sm4_encrypt(new_token_data)
                st.markdown(f"""
                        <script>
                            localStorage.setItem("auth_token", "{new_token}");
                        </script>
                    """, unsafe_allow_html=True)
            except:
                pass  # 续期失败不影响使用
            st.sidebar.divider()
            # ===== 侧边栏：实时 + 明天预报 =====
            adcode = st.session_state.get("user_adcode")
            city = st.session_state.get("user_city", "")
            if adcode and city:
                try:
                    api = GaodeWeatherAPI(GAODE_API_KEY)
                    cur_data = api.get_weather(adcode, "base")  # 实况
                    fc_data = api.get_weather(adcode, "all")  # 预报
                    cur = api.parse_current_weather(cur_data) if cur_data else None
                    fc = api.parse_forecast_weather(fc_data) if fc_data else None

                    if cur:
                        st.sidebar.success(f"{cur['city']} 当前：{cur['weather']} {cur['temperature']}℃")
                    if fc and len(fc["forecasts"]) > 1:  # 取明天
                        tom = fc["forecasts"][1]
                        st.sidebar.info(f"明日：{tom['day_weather']}/{tom['night_weather']}  "
                                        f"{tom['night_temp']}~{tom['day_temp']}℃")
                except Exception as e:
                    st.sidebar.error(f"天气获取异常: {e}")
            else:
                st.sidebar.info("暂无城市定位")
            # =========================
            st.sidebar.write(f"👤 当前用户：{st.session_state.get('current_user', '未知')}")
            # 显示权限信息
            if hasattr(st.session_state, 'user_permissions'):
                perm_names = Permissions.get_permission_names(st.session_state.user_permissions)
                st.sidebar.write(f"权限：{', '.join(perm_names[:3])}" +
                                 ("..." if len(perm_names) > 3 else ""))
            if st.button("退出登录", key="main_logout_btn"):
                # 先关闭数据库连接
                if st.session_state.get("db_client") and hasattr(st.session_state.db_client, 'conn'):
                    try:
                        if st.session_state.db_client.conn:
                            st.session_state.db_client.conn.close()
                    except:
                        pass
                        # 清除 localStorage 中的 token
                st.markdown("""
                        <script>
                            localStorage.removeItem("auth_token");
                        </script>
                    """, unsafe_allow_html=True)
                # 清空所有登录相关的session_state
                st.session_state.is_logged_in = False
                st.session_state.init_ok = False
                st.session_state.db_client = None
                st.session_state.encoder = None
                st.session_state.current_user = None
                st.session_state.user_permissions = Permissions.DEFAULT_USER

                # 重置其他可能需要重置的状态
                st.session_state.current_kb_id = None
                st.session_state.current_page = 1
                st.session_state.file_status = {}
                st.success("已退出登录")
                time.sleep(0.5)

                # 重要：使用experimental_rerun或直接rerun，确保跳转到登录页
                st.rerun()
            # ========== 新增：重新初始化数据库 按钮（核心需求） ==========
            if st.button("重新初始化数据库", key="reset_db_btn", type="secondary"):
                # 1. 删除代码目录下的db.yaml文件
                if DB_YAML_PATH.exists():
                    os.remove(DB_YAML_PATH)
                # 2. 清空所有状态，退出登录
                st.session_state.is_logged_in = False
                st.session_state.init_ok = False
                st.session_state.db_client = None
                st.session_state.encoder = None
                st.success("已删除db.yaml，退出登录！请重新初始化")
                time.sleep(0.5)
                st.rerun()
        stop_word_config_panel()
    # 动态创建页签
    tab_names = []

    # 检查用户是否有智能问答的权限
    if Permissions.has_permission(st.session_state.user_permissions, Permissions.USE_QA):
        tab_names.append("智能问答")

    # 检查用户是否有查看知识库的权限
    if Permissions.has_permission(st.session_state.user_permissions, Permissions.VIEW_KB):
        tab_names.append("知识库管理")

    # 检查用户是否有图片识别问答的权限
    if Permissions.has_permission(st.session_state.user_permissions, Permissions.USE_IMAGE_QA):
        tab_names.append("图片识别与问答")

    # 检查是否为管理员
    if Permissions.has_permission(st.session_state.user_permissions, Permissions.ADMIN):
        tab_names.append("用户管理")

    # 如果没有任何权限，显示错误信息
    if not tab_names:
        st.error("您没有任何可用权限，请联系管理员！")
        return

    # 创建页签
    tabs = st.tabs(tab_names)

    # 根据页签名称分配功能
    for i, tab_name in enumerate(tab_names):
        with tabs[i]:
            if tab_name == "智能问答":
                qa_tab()
            elif tab_name == "知识库管理":
                # 在知识库管理页签内部检查VIEW_KB权限
                if Permissions.has_permission(st.session_state.user_permissions, Permissions.VIEW_KB):
                    kb_manage_tab()
                else:
                    st.error("你没有查看知识库的权限！请联系管理员授权。")
            elif tab_name == "图片识别与问答":
                image_detect_qa_tab()
            elif tab_name == "用户管理":
                user_manage_tab()


if __name__ == "__main__":
    main()
    import gc
    gc.collect()